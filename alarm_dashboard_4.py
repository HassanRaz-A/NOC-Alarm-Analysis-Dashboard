import io
import inspect as _inspect
from datetime import datetime, timedelta

import numpy as np
import pandas as pd
import streamlit as st
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN

# ── Optional ML stack (graceful if scikit-learn is not installed) ─────────────
try:
    import joblib
    from sklearn.ensemble import GradientBoostingRegressor
    from sklearn.metrics import mean_absolute_error, mean_squared_error, r2_score
    SKLEARN_AVAILABLE = True
except ImportError:
    SKLEARN_AVAILABLE = False

# ── xlrd version compatibility ───────────────────────────────────────────────
try:
    import xlrd
    XLRD_AVAILABLE = True
    XLRD_SUPPORTS_CORRUPTION_FLAG = (
        "ignore_workbook_corruption" in _inspect.signature(xlrd.open_workbook).parameters
    )
except ImportError:
    XLRD_AVAILABLE = False
    XLRD_SUPPORTS_CORRUPTION_FLAG = False

# ── Auto-install python-calamine (handles corrupted .xls files that xlrd 2.x cannot) ──
try:
    import python_calamine  # noqa: F401 — just checking it is available
    CALAMINE_AVAILABLE = True
except ImportError:
    try:
        import subprocess, sys
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "python-calamine", "-q"],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        CALAMINE_AVAILABLE = True
    except Exception:
        CALAMINE_AVAILABLE = False

# ── Severity weights ─────────────────────────────────────────────────────────
SEVERITY_MAP = {
    "critical": 4, "major": 3, "minor": 2,
    "warning": 1,  "warn": 1,  "informational": 0, "info": 0,
}
LEVEL_COLOR = {
    "Critical": "#e53935", "Major": "#fb8c00",
    "Minor":    "#fdd835", "Warning": "#43a047",
}

# ════════════════════════════════════════════════════════════════════════════
# PAGE CONFIG + THEME
# ════════════════════════════════════════════════════════════════════════════
st.set_page_config(page_title="NOC Alarm Intelligence", layout="wide", page_icon="🔔")

if "theme_mode" not in st.session_state:
    st.session_state.theme_mode = "Dark"

theme_mode = st.sidebar.selectbox(
    "🎨 Theme", ["Dark", "Light"],
    index=0 if st.session_state.theme_mode == "Dark" else 1,
)
st.session_state.theme_mode = theme_mode

if theme_mode == "Dark":
    bg = "#0d1117"; card = "#161b22"; border = "#30363d"
    text = "#c9d1d9"; muted = "#8b949e"
    blue = "#1e90ff"; blue_s = "rgba(30,144,255,0.35)"
    purple_s = "rgba(137,0,255,0.22)"; tmpl = "plotly_dark"
else:
    bg = "#f8fafc"; card = "#ffffff"; border = "#d0d7de"
    text = "#111827"; muted = "#4b5563"
    blue = "#2563eb"; blue_s = "rgba(37,99,235,0.18)"
    purple_s = "rgba(124,58,237,0.14)"; tmpl = "plotly_white"

st.markdown(f"""<style>
:root{{--bg:{bg};--card:{card};--border:{border};--text:{text};--blue:{blue};}}
.stApp{{background:radial-gradient(circle at top right,{purple_s},transparent 22%),
        radial-gradient(circle at top left,{blue_s},transparent 20%),{bg};color:{text};}}
html,body,[class*="css"]{{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif;}}
h1,h2,h3,h4,h5,h6,p,div,span,label{{color:{text} !important;}}
.block-container{{padding-top:1rem;padding-bottom:1.2rem;}}
[data-testid="stSidebar"]{{background:linear-gradient(180deg,rgba(30,144,255,.05),rgba(137,0,255,.03)),{card};border-right:1px solid {border};}}
div[data-testid="stMetric"]{{background:linear-gradient(180deg,rgba(30,144,255,.04),rgba(137,0,255,.02)),{card};border:1px solid {border};border-radius:14px;padding:14px;transition:all .28s ease;}}
div[data-testid="stMetric"]:hover{{border-color:{blue};box-shadow:0 0 18px {blue_s};transform:translateY(-2px);}}
div[data-testid="stDataFrame"],div[data-testid="stTable"]{{background:{card};border:1px solid {border};border-radius:14px;padding:6px;}}
div[data-testid="stDataFrame"]:hover,div[data-testid="stTable"]:hover{{border-color:{blue};box-shadow:0 0 16px {blue_s};}}
.stButton>button,.stDownloadButton>button{{background:transparent!important;color:{blue}!important;border:2px solid {blue}!important;border-radius:10px!important;font-weight:700!important;transition:all .22s!important;}}
.stButton>button:hover,.stDownloadButton>button:hover{{background:{blue}!important;color:white!important;box-shadow:0 0 20px {blue}!important;transform:translateY(-1px);}}
div[data-baseweb="select"]>div{{background:{card}!important;border:1px solid {border}!important;border-radius:10px!important;}}
[data-testid="stFileUploader"]{{background:linear-gradient(180deg,rgba(30,144,255,.03),rgba(137,0,255,.02)),{card};border:1px solid {border};border-radius:14px;padding:10px;}}
h1{{text-shadow:0 0 18px {blue_s};}}

/* ── App header banner ─────────────────────────────────────────────── */
.app-header{{position:relative;border:1px solid {border};border-radius:18px;padding:18px 24px;margin-bottom:14px;
  background:linear-gradient(120deg,rgba(30,144,255,.10),rgba(137,0,255,.06) 60%,transparent),{card};
  box-shadow:0 0 26px {blue_s};overflow:hidden;}}
.app-header::after{{content:"";position:absolute;inset:0;background:radial-gradient(circle at 90% -20%,{blue_s},transparent 55%);pointer-events:none;}}
.app-header .ah-title{{font-size:1.9rem;font-weight:800;letter-spacing:.3px;
  background:linear-gradient(90deg,{blue},#8a5cff 70%);-webkit-background-clip:text;background-clip:text;-webkit-text-fill-color:transparent;}}
.app-header .ah-sub{{color:{muted}!important;font-size:.95rem;margin-top:2px;}}

/* ── Context bar (active file + KPIs) ──────────────────────────────── */
.context-bar{{display:flex;flex-wrap:wrap;gap:10px;align-items:stretch;margin:4px 0 16px;}}
.ctx-chip{{flex:1 1 150px;border:1px solid {border};border-radius:12px;padding:10px 14px;background:{card};
  transition:all .25s ease;}}
.ctx-chip:hover{{border-color:{blue};box-shadow:0 0 14px {blue_s};transform:translateY(-2px);}}
.ctx-chip .cl{{color:{muted}!important;font-size:.72rem;text-transform:uppercase;letter-spacing:.6px;}}
.ctx-chip .cv{{font-size:1.15rem;font-weight:700;}}
.ctx-active{{border-left:3px solid {blue};}}

/* ── Hero (landing) ────────────────────────────────────────────────── */
.hero{{text-align:center;padding:44px 18px 26px;border:1px solid {border};border-radius:22px;margin-bottom:22px;
  background:radial-gradient(circle at 50% -30%,{blue_s},transparent 60%),{card};box-shadow:0 0 34px {blue_s};}}
.hero .h-badge{{display:inline-block;border:1px solid {blue};color:{blue}!important;border-radius:999px;
  padding:4px 14px;font-size:.78rem;font-weight:700;letter-spacing:1px;text-transform:uppercase;margin-bottom:14px;}}
.hero h1.h-title{{font-size:2.9rem;line-height:1.1;font-weight:900;margin:0;
  background:linear-gradient(90deg,{blue},#8a5cff 75%);-webkit-background-clip:text;background-clip:text;-webkit-text-fill-color:transparent;text-shadow:none;}}
.hero .h-tag{{color:{muted}!important;font-size:1.08rem;margin:14px auto 0;max-width:680px;}}
.h-pills{{display:flex;flex-wrap:wrap;gap:8px;justify-content:center;margin-top:18px;}}
.h-pill{{border:1px solid {border};border-radius:999px;padding:6px 14px;font-size:.82rem;background:{bg};}}
.h-pill:hover{{border-color:{blue};box-shadow:0 0 12px {blue_s};}}

/* ── Feature cards ─────────────────────────────────────────────────── */
.feature-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(230px,1fr));gap:14px;margin:6px 0 26px;}}
.feature-card{{border:1px solid {border};border-radius:16px;padding:18px;background:{card};transition:all .28s ease;}}
.feature-card:hover{{border-color:{blue};box-shadow:0 0 20px {blue_s};transform:translateY(-3px);}}
.feature-card .fc-icon{{font-size:1.7rem;}}
.feature-card .fc-title{{font-weight:800;font-size:1.05rem;margin:6px 0 4px;}}
.feature-card .fc-desc{{color:{muted}!important;font-size:.88rem;line-height:1.4;}}

/* ── Quick start ───────────────────────────────────────────────────── */
.quickstart{{border:1px solid {border};border-radius:18px;padding:20px 24px;background:{card};margin-bottom:18px;}}
.quickstart .qs-h{{font-weight:800;font-size:1.15rem;margin-bottom:12px;}}
.qs-steps{{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:14px;}}
.qs-step{{display:flex;gap:12px;align-items:flex-start;}}
.qs-num{{flex:0 0 30px;height:30px;width:30px;border-radius:50%;display:flex;align-items:center;justify-content:center;
  font-weight:800;background:linear-gradient(135deg,{blue},#8a5cff);color:#fff!important;box-shadow:0 0 12px {blue_s};}}
.qs-txt{{font-size:.9rem;color:{text}!important;}}
.qs-txt b{{color:{text}!important;}}
.cta{{text-align:center;border:1px dashed {blue};border-radius:16px;padding:16px;margin-top:8px;background:{bg};font-weight:600;}}

/* ── Section title ─────────────────────────────────────────────────── */
.section-title{{font-size:1.15rem;font-weight:800;margin:6px 0 2px;padding-left:11px;border-left:4px solid {blue};}}
</style>""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════════════════
# UTILITY FUNCTIONS
# ════════════════════════════════════════════════════════════════════════════

def short_device(full: str) -> str:
    """Last segment of a dash-separated device path.
    'ODU#1-DOU2-eHUB#1-eRDU2'  ->  'eRDU2'
    """
    parts = str(full).strip().split("-")
    return parts[-1].strip() if parts else str(full)


def apply_theme(fig):
    fig.update_layout(
        template=tmpl,
        paper_bgcolor=bg,
        plot_bgcolor=card if theme_mode == "Dark" else bg,
        font=dict(color=text),
        title_font=dict(color=text, size=18),
        legend=dict(font=dict(color=text)),
        margin=dict(l=40, r=30, t=55, b=40),
    )
    return fig


def fig_to_png(fig):
    return pio.to_image(fig, format="png", width=1400, height=800, scale=2)


def dl_btn(fig, fname, label):
    try:
        st.download_button(label, fig_to_png(fig), fname, "image/png", key=fname)
    except Exception:
        st.info("PNG export requires kaleido:  pip install kaleido")


PCFG = {
    "displaylogo": False,
    "toImageButtonOptions": {"format":"png","filename":"chart","height":700,"width":1200,"scale":2},
}


# ════════════════════════════════════════════════════════════════════════════
# XLRD COMPAT
# ════════════════════════════════════════════════════════════════════════════

def _open_xlrd(file_contents: bytes):
    """Works with xlrd 1.x (has corruption flag) and xlrd 2.x (removed it)."""
    if XLRD_SUPPORTS_CORRUPTION_FLAG:
        try:
            return xlrd.open_workbook(file_contents=file_contents, ignore_workbook_corruption=True)
        except Exception:
            pass
    return xlrd.open_workbook(file_contents=file_contents)


# ════════════════════════════════════════════════════════════════════════════
# FILE LOADING  (cached)
# ════════════════════════════════════════════════════════════════════════════

@st.cache_data(show_spinner=False)
def get_sheets(file_bytes: bytes, file_name: str) -> list:
    name = file_name.lower()
    if name.endswith(".xlsx"):
        return pd.ExcelFile(io.BytesIO(file_bytes), engine="openpyxl").sheet_names
    if name.endswith(".xls"):
        # Try calamine first — it handles many corrupted .xls files better than xlrd 2.x
        try:
            return pd.ExcelFile(io.BytesIO(file_bytes), engine="calamine").sheet_names
        except Exception:
            pass
        if not XLRD_AVAILABLE:
            raise ImportError(
                "Could not read this .xls file. Try: pip install xlrd==1.2.0"
            )
        try:
            return pd.ExcelFile(io.BytesIO(file_bytes), engine="xlrd").sheet_names
        except Exception:
            try:
                return _open_xlrd(file_bytes).sheet_names()
            except Exception:
                raise ValueError(
                    "Corrupted .xls file. Fix: pip install xlrd==1.2.0  "
                    "(xlrd 2.x removed support for corrupted workbooks)"
                )
    raise ValueError("Unsupported format — upload .xls or .xlsx")


def _unwrap_dict(result, sheet: str) -> pd.DataFrame:
    """If pd.read_excel returns a dict (e.g. calamine quirk), extract the right sheet."""
    if isinstance(result, dict):
        if sheet in result:
            return result[sheet]
        if result:
            return list(result.values())[0]
        raise ValueError("No sheets found in workbook.")
    return result


@st.cache_data(show_spinner=False)
def read_sheet(file_bytes: bytes, file_name: str, sheet: str) -> pd.DataFrame:
    name = file_name.lower()

    # calamine is the fastest engine and handles many corrupted files
    try:
        result = pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet, engine="calamine")
        return _unwrap_dict(result, sheet)
    except Exception:
        pass

    if name.endswith(".xlsx"):
        result = pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet, engine="openpyxl")
        return _unwrap_dict(result, sheet)

    if name.endswith(".xls"):
        if not XLRD_AVAILABLE:
            raise ImportError(
                "Could not read this .xls file. Try: pip install xlrd==1.2.0"
            )
        try:
            result = pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet, engine="xlrd")
            return _unwrap_dict(result, sheet)
        except Exception:
            pass
        # Last resort: open with xlrd directly (handles some corruption)
        try:
            book = _open_xlrd(file_bytes)
            s = book.sheet_by_name(sheet)
            data = [s.row_values(r) for r in range(s.nrows)]
            if not data:
                raise ValueError("Sheet is empty")
            return pd.DataFrame(data[1:], columns=[str(h).strip() for h in data[0]])
        except Exception:
            raise ValueError(
                "Corrupted .xls file could not be read. "
                "Fix: pip install xlrd==1.2.0"
            )

    raise ValueError("Unsupported format")


# ════════════════════════════════════════════════════════════════════════════
# DATA CLEANING
# ════════════════════════════════════════════════════════════════════════════

def clean_df(df) -> pd.DataFrame:
    # Guard: if somehow a dict slips through, take first sheet
    if isinstance(df, dict):
        df = list(df.values())[0] if df else pd.DataFrame()
    df = df.copy()
    df.columns = df.columns.astype(str).str.strip()

    for col in ["No.", "Type", "Note"]:          # drop useless columns
        if col in df.columns:
            df.drop(columns=col, inplace=True)

    for col in ["Time Received", "Time Cleared"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce")

    if "Time Received" in df.columns and "Time Cleared" in df.columns:
        df["Duration (Min)"] = (
            (df["Time Cleared"] - df["Time Received"]).dt.total_seconds() / 60
        ).clip(lower=0)
    else:
        df["Duration (Min)"] = None

    for col in ["Level", "BIU", "Device", "Name", "Message"]:
        if col in df.columns:
            df[col] = df[col].fillna("Unknown").astype(str).str.strip()

    if "Device" in df.columns:
        df["Device Short"] = df["Device"].apply(short_device)

    if "Level" in df.columns:
        df["Severity"] = df["Level"].str.lower().map(SEVERITY_MAP).fillna(1).astype(float)

    if "Time Received" in df.columns:
        df["Hour"]      = df["Time Received"].dt.hour
        df["DayOfWeek"] = df["Time Received"].dt.day_name()
        df["Date"]      = df["Time Received"].dt.date

    return df


# ════════════════════════════════════════════════════════════════════════════
# DEVICE HEALTH SCORING + PREDICTION ENGINE
# ════════════════════════════════════════════════════════════════════════════

def compute_health(df: pd.DataFrame) -> pd.DataFrame:
    """
    Returns per-device health metrics:
      Health Score 0-100  (100 = healthy, 0 = critical)
      MTBF  — Mean Time Between Failures (hours)
      MTTR  — Mean Time To Repair (minutes = avg Duration)
      Predicted Next Alarm + Confidence %
    """
    if "Device" not in df.columns or df.empty:
        return pd.DataFrame()

    span_days = max(
        (df["Time Received"].max() - df["Time Received"].min()).total_seconds() / 86400, 1
    )

    records = []
    for device, grp in df.groupby("Device"):
        grp   = grp.sort_values("Time Received")
        n     = len(grp)
        times = grp["Time Received"].dropna()

        avg_sev = float(grp["Severity"].mean()) if "Severity" in grp.columns else 1.0
        avg_dur = grp["Duration (Min)"].dropna().mean() if "Duration (Min)" in grp.columns else 0
        avg_dur = float(avg_dur) if not pd.isna(avg_dur) else 0.0
        alarms_per_day = n / span_days

        # ── MTBF + exponentially-weighted prediction ──────────────────────
        mtbf_h = predicted = conf = None
        if len(times) >= 2:
            diffs_h = times.diff().dropna().dt.total_seconds() / 3600
            mtbf_h  = round(float(diffs_h.mean()), 2)
            ewm_iat = float(diffs_h.ewm(span=min(3, len(diffs_h))).mean().iloc[-1])
            std_iat = float(diffs_h.std()) if len(diffs_h) > 1 else ewm_iat
            cv      = std_iat / (ewm_iat + 1e-9)
            predicted = times.iloc[-1] + timedelta(hours=max(ewm_iat, 0.1))
            conf    = round(max(0.0, min(100.0, 100.0 * (1.0 - cv))), 1)

        top_name  = grp["Name"].value_counts().idxmax()  if "Name"  in grp.columns else "N/A"
        top_level = grp["Level"].value_counts().idxmax() if "Level" in grp.columns else "N/A"
        top_biu   = grp["BIU"].value_counts().idxmax()   if "BIU"   in grp.columns else "N/A"

        records.append({
            "Device Full":    device,
            "Device":         short_device(device),
            "BIU":            top_biu,
            "Alarms":         n,
            "Alarms/Day":     round(alarms_per_day, 2),
            "Avg Severity":   round(avg_sev, 2),
            "MTTR (Min)":     round(avg_dur, 1),
            "MTBF (Hrs)":     mtbf_h,
            "Last Alarm":     times.iloc[-1] if len(times) > 0 else None,
            "Predicted Next": predicted,
            "Confidence (%)": conf,
            "Top Level":      top_level,
            "Top Alarm":      top_name,
        })

    hdf = pd.DataFrame(records)
    if hdf.empty:
        return hdf

    def norm(s: pd.Series) -> pd.Series:
        s = s.fillna(s.median())
        mn, mx = s.min(), s.max()
        return (s - mn) / (mx - mn) if mx > mn else pd.Series(0.5, index=s.index)

    freq_n  = norm(hdf["Alarms/Day"])
    sev_n   = norm(hdf["Avg Severity"])
    dur_n   = norm(hdf["MTTR (Min)"])
    mtbf_n  = 1 - norm(hdf["MTBF (Hrs)"].fillna(hdf["MTBF (Hrs)"].max()))

    penalty = 0.35*freq_n + 0.35*sev_n + 0.15*dur_n + 0.15*mtbf_n
    hdf["Health Score"] = (100.0*(1.0 - penalty)).clip(0, 100).round(1)

    def status(s):
        if s >= 80: return "🟢 Good"
        if s >= 60: return "🟡 Fair"
        if s >= 40: return "🟠 Poor"
        return "🔴 Critical"

    hdf["Status"] = hdf["Health Score"].apply(status)
    return hdf.sort_values("Health Score").reset_index(drop=True)


# ════════════════════════════════════════════════════════════════════════════
# ML FORECASTING  (trainable daily alarm-volume model)
# ════════════════════════════════════════════════════════════════════════════

ML_FEATURES = ["dow", "dom", "month", "is_weekend", "lag1", "lag7", "roll3", "roll7"]


def build_daily_series(df: pd.DataFrame) -> pd.DataFrame:
    """Continuous daily alarm-count series (calendar gaps filled with 0)."""
    if "Time Received" not in df.columns:
        return pd.DataFrame(columns=["Date", "Count"])
    tr = df["Time Received"].dropna()
    if tr.empty:
        return pd.DataFrame(columns=["Date", "Count"])
    daily = tr.dt.floor("D").value_counts().sort_index()
    full = pd.date_range(daily.index.min(), daily.index.max(), freq="D")
    daily = daily.reindex(full, fill_value=0)
    out = daily.reset_index()
    out.columns = ["Date", "Count"]
    return out


def _add_ml_features(daily: pd.DataFrame) -> pd.DataFrame:
    d = daily.copy()
    d["Date"]       = pd.to_datetime(d["Date"])
    d["dow"]        = d["Date"].dt.dayofweek
    d["dom"]        = d["Date"].dt.day
    d["month"]      = d["Date"].dt.month
    d["is_weekend"] = (d["dow"] >= 5).astype(int)
    d["lag1"]       = d["Count"].shift(1)
    d["lag7"]       = d["Count"].shift(7)
    d["roll3"]      = d["Count"].shift(1).rolling(3).mean()
    d["roll7"]      = d["Count"].shift(1).rolling(7).mean()
    return d


@st.cache_data(show_spinner=False)
def train_volume_model(daily: pd.DataFrame):
    """Train a daily-volume forecaster. Returns a dict with model + holdout metrics."""
    if not SKLEARN_AVAILABLE:
        return {"ok": False, "reason": "scikit-learn not installed"}

    feat = _add_ml_features(daily).dropna().reset_index(drop=True)
    if len(feat) < 10:
        return {"ok": False, "reason": f"only {len(feat)} usable training rows (need ≥ 10)"}

    X, y = feat[ML_FEATURES], feat["Count"]
    n = len(feat)
    split = max(int(n * 0.8), n - 14)          # last ≤14 rows held out
    split = min(split, n - 2)                  # guarantee ≥2 holdout rows

    def _fit(Xt, yt):
        m = GradientBoostingRegressor(
            n_estimators=200, max_depth=3, learning_rate=0.05, random_state=42
        )
        m.fit(Xt, yt)
        return m

    eval_model = _fit(X.iloc[:split], y.iloc[:split])
    y_pred = eval_model.predict(X.iloc[split:]).clip(min=0)
    y_true = y.iloc[split:].values

    metrics = {
        "MAE":  float(mean_absolute_error(y_true, y_pred)),
        "RMSE": float(np.sqrt(mean_squared_error(y_true, y_pred))),
        "R2":   float(r2_score(y_true, y_pred)) if len(y_true) > 1 else float("nan"),
    }
    eval_df = pd.DataFrame({
        "Date":      feat["Date"].iloc[split:].values,
        "Actual":    y_true,
        "Predicted": np.round(y_pred, 1),
    })

    full_model = _fit(X, y)                    # refit on everything for forecasting
    return {
        "ok": True, "model": full_model, "metrics": metrics,
        "eval_df": eval_df, "n_train": split, "n_test": n - split,
        "importances": dict(zip(ML_FEATURES, full_model.feature_importances_)),
    }


def forecast_volume(model, daily: pd.DataFrame, horizon: int) -> pd.DataFrame:
    """Recursive multi-step forecast of daily alarm counts."""
    hist = _add_ml_features(daily)[["Date", "Count"]].copy()
    rows = []
    last_known = float(hist["Count"].iloc[-1])
    for _ in range(horizon):
        nxt = hist["Date"].iloc[-1] + pd.Timedelta(days=1)
        tmp = pd.concat([hist, pd.DataFrame([{"Date": nxt, "Count": np.nan}])],
                        ignore_index=True)
        # 1-row DataFrame (keeps feature names) of floats (avoids object-dtype fillna).
        feat_row = _add_ml_features(tmp).iloc[[-1]][ML_FEATURES].astype(float)
        # roll/lag features can be NaN on very short history → fall back to last count.
        feat_row = feat_row.fillna(last_known)
        yhat = max(0.0, float(model.predict(feat_row)[0]))
        rows.append({"Date": nxt, "Forecast": round(yhat, 1)})
        hist = pd.concat([hist, pd.DataFrame([{"Date": nxt, "Count": yhat}])],
                         ignore_index=True)
        last_known = yhat
    return pd.DataFrame(rows)


def load_alarm_history(file_bytes: bytes, file_name: str) -> pd.DataFrame:
    """Load + clean the alarm sheet of an uploaded workbook (auto-picks EventLog)."""
    sheets = get_sheets(file_bytes, file_name)
    if not sheets:
        return pd.DataFrame()
    pick = next((s for s in sheets if str(s).strip().lower() == "eventlog"), sheets[0])
    return clean_df(read_sheet(file_bytes, file_name, pick))


def combine_histories(frames: list) -> pd.DataFrame:
    """Concatenate alarm frames and drop duplicate alarms (overlapping exports)."""
    frames = [f for f in frames if f is not None and not f.empty]
    if not frames:
        return pd.DataFrame()
    combined = pd.concat(frames, ignore_index=True)
    keys = [c for c in ["Time Received", "Device", "Name", "Message", "Level", "BIU"]
            if c in combined.columns]
    if keys:
        combined = combined.drop_duplicates(subset=keys).reset_index(drop=True)
    return combined


# ════════════════════════════════════════════════════════════════════════════
# PPT EXPORT
# ════════════════════════════════════════════════════════════════════════════

def _summary_sections(df):
    total = len(df)
    ul = df["Level"].nunique()  if "Level"  in df.columns else 0
    ub = df["BIU"].nunique()    if "BIU"    in df.columns else 0
    ud = df["Device"].nunique() if "Device" in df.columns else 0
    ad = round(df["Duration (Min)"].dropna().mean(), 1) if "Duration (Min)" in df.columns else 0
    tl = df["Level"].value_counts().idxmax()  if "Level"  in df.columns else "N/A"
    tb = df["BIU"].value_counts().idxmax()    if "BIU"    in df.columns else "N/A"
    td = short_device(df["Device"].value_counts().idxmax()) if "Device" in df.columns else "N/A"
    return [
        ("Dataset Overview", [
            f"Total alarms: {total}",
            f"Unique levels: {ul}  |  Unique BIUs: {ub}  |  Unique devices: {ud}",
        ]),
        ("Key Highlights", [
            f"Most frequent level: {tl}",
            f"Most affected BIU: {tb}",
            f"Most affected device: {td}",
            f"Avg alarm duration: {ad} minutes",
        ]),
    ]


def _write_tf(tf, sections):
    tf.clear()
    first = True
    for heading, bullets in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = heading; p.level = 0; p.font.bold = True; p.font.size = Pt(20)
        for b in bullets:
            bp = tf.add_paragraph(); bp.text = b; bp.level = 1; bp.font.size = Pt(16)


def build_ppt(df, source_file, sheet_name, template_bytes=None):
    prs      = Presentation(io.BytesIO(template_bytes)) if template_bytes else Presentation()
    sections = _summary_sections(df)

    if not template_bytes:
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        try:
            if sl.shapes.title: sl.shapes.title.text = "Alarm Analysis Report"
        except Exception: pass
        try: sl.placeholders[1].text = f"File: {source_file} | Sheet: {sheet_name}"
        except Exception: pass
        sl2 = prs.slides.add_slide(prs.slide_layouts[1])
        try:
            if sl2.shapes.title: sl2.shapes.title.text = "Executive Summary"
        except Exception: pass
        try: _write_tf(sl2.shapes.placeholders[1].text_frame, sections)
        except Exception: pass

    def add_chart(title, cats, vals, ctype=XL_CHART_TYPE.COLUMN_CLUSTERED):
        sl = prs.slides.add_slide(prs.slide_layouts[5])
        try:
            if sl.shapes.title: sl.shapes.title.text = title
        except Exception: pass
        cd = CategoryChartData()
        cd.categories = [str(c) for c in cats]
        cd.add_series("Count", [int(v) for v in vals])
        gf = sl.shapes.add_chart(ctype, 700000, 1200000, 7900000, 4700000, cd)
        gf.chart.has_legend = True
        gf.chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    if "Level" in df.columns and not df.empty:
        lx = df["Level"].value_counts().reset_index(); lx.columns = ["Level","Count"]
        add_chart("Alarm Count by Level", lx["Level"], lx["Count"])
        add_chart("Level Distribution",   lx["Level"], lx["Count"], XL_CHART_TYPE.PIE)
    if "BIU" in df.columns and not df.empty:
        bx = df["BIU"].value_counts().reset_index().head(15); bx.columns = ["BIU","Count"]
        add_chart("Alarm Count by BIU", bx["BIU"], bx["Count"])
    if "Device" in df.columns and not df.empty:
        dx = df["Device"].value_counts().reset_index().head(15); dx.columns = ["Device","Count"]
        dx["Device"] = dx["Device"].apply(short_device)
        add_chart("Top Affected Devices", dx["Device"], dx["Count"])
    if "Name" in df.columns and not df.empty:
        nx = df["Name"].value_counts().reset_index().head(12); nx.columns = ["Name","Count"]
        add_chart("Top Alarm Names", nx["Name"], nx["Count"], XL_CHART_TYPE.BAR_CLUSTERED)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf


# ════════════════════════════════════════════════════════════════════════════
# MAIN UI
# ════════════════════════════════════════════════════════════════════════════

def render_landing():
    """Pre-upload landing page: hero + feature cards + quick-start + CTA."""
    st.markdown(
        """
        <div class="hero">
          <div class="h-badge">DTS Wireless · NOC</div>
          <h1 class="h-title">NOC Alarm Intelligence</h1>
          <div class="h-tag">Turn raw alarm EventLog exports into health scores, predictive
          insight, peak-hour patterns and editable PowerPoint reports — all in your browser.</div>
          <div class="h-pills">
            <span class="h-pill">📊 Analytics</span>
            <span class="h-pill">🏥 Device Health</span>
            <span class="h-pill">🔮 Predictions</span>
            <span class="h-pill">🤖 ML Forecast</span>
            <span class="h-pill">⚖️ File Comparison</span>
            <span class="h-pill">📥 PPT / CSV Export</span>
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    features = [
        ("📊", "Overview", "KPIs, alarms by level &amp; BIU, and a daily timeline at a glance."),
        ("🔍", "Alarm Analysis", "Level pies, duration boxes, device &amp; BIU heatmaps, sunburst hierarchy."),
        ("🏥", "Device Health", "0–100 health score per device with MTBF, MTTR and a reliability map."),
        ("🕐", "Time Patterns", "Hour-of-day, day-of-week and peak-hour heatmaps to spot recurring spikes."),
        ("🔮", "Predictions", "Per-device next-alarm estimate &amp; urgency from inter-arrival statistics."),
        ("🤖", "ML Forecast", "A trained model that forecasts upcoming daily alarm volume."),
    ]
    cards = "".join(
        f'<div class="feature-card"><div class="fc-icon">{i}</div>'
        f'<div class="fc-title">{t}</div><div class="fc-desc">{d}</div></div>'
        for i, t, d in features
    )
    st.markdown(f'<div class="feature-grid">{cards}</div>', unsafe_allow_html=True)

    st.markdown(
        """
        <div class="quickstart">
          <div class="qs-h">🚀 Quick Start</div>
          <div class="qs-steps">
            <div class="qs-step"><div class="qs-num">1</div><div class="qs-txt">
              Export the <b>EventLog</b> alarm report (<b>.xls</b> or <b>.xlsx</b>) from your alarm system.</div></div>
            <div class="qs-step"><div class="qs-num">2</div><div class="qs-txt">
              In the <b>sidebar</b>, upload one or more files under <b>📂 Data Sources</b>.</div></div>
            <div class="qs-step"><div class="qs-num">3</div><div class="qs-txt">
              Pick the <b>workbook</b> and <b>sheet</b> to analyse (multi-file enables Compare mode).</div></div>
            <div class="qs-step"><div class="qs-num">4</div><div class="qs-txt">
              Explore the tabs, then download a <b>PowerPoint</b> report or <b>CSV</b> tables.</div></div>
          </div>
          <div class="cta">⬅️ &nbsp; Upload your Excel files in the sidebar to begin</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


st.markdown(
    """
    <div class="app-header">
      <div class="ah-title">🔔 NOC Alarm Intelligence Dashboard</div>
      <div class="ah-sub">Multi-file upload · Alarm analysis · Device health scoring ·
      Predictive analytics · ML forecasting · File comparison</div>
    </div>
    """,
    unsafe_allow_html=True,
)

# ── Sidebar: upload ──────────────────────────────────────────────────────────
st.sidebar.header("📂 Data Sources")
uploaded = st.file_uploader("Upload Excel Files", type=["xls","xlsx"], accept_multiple_files=True)
ppt_tpl  = st.file_uploader("Optional: PowerPoint Template", type=["pptx"])

wmap = {}
if uploaded:
    for uf in uploaded:
        try:
            wmap[uf.name] = {"file": uf, "sheets": get_sheets(uf.getvalue(), uf.name)}
        except Exception as e:
            st.sidebar.error(f"❌ {uf.name}: {e}")

if not wmap:
    render_landing()
    st.stop()

# ── Sidebar: workbook / sheet selector ──────────────────────────────────────
st.sidebar.header("📋 Active Sheet")
wb_list = list(wmap.keys())
sel_wb  = st.sidebar.selectbox("Workbook", wb_list)

# Guard: selectbox returns None before Streamlit fully initialises
# (also triggered when running via python.exe instead of streamlit run)
if not sel_wb or sel_wb not in wmap:
    st.info("⬆️  Upload one or more Excel files and use  streamlit run  to launch.")
    st.stop()

sheet_list = wmap[sel_wb]["sheets"]
sel_sheet  = st.sidebar.selectbox("Sheet", sheet_list)

if not sel_sheet:
    st.warning("No sheets found in this workbook.")
    st.stop()

uf_obj = wmap[sel_wb]["file"]
raw    = read_sheet(uf_obj.getvalue(), uf_obj.name, sel_sheet)
df     = clean_df(raw)

# ── Sidebar: filters ─────────────────────────────────────────────────────────
st.sidebar.header("🔎 Filters")
filtered = df.copy()

if "Level" in filtered.columns:
    opts = sorted(filtered["Level"].dropna().unique())
    sel  = st.sidebar.multiselect("Level", opts, default=list(opts))
    filtered = filtered[filtered["Level"].isin(sel)]

if "BIU" in filtered.columns:
    opts = sorted(filtered["BIU"].dropna().unique())
    sel  = st.sidebar.multiselect("BIU", opts, default=list(opts))
    filtered = filtered[filtered["BIU"].isin(sel)]

if "Time Received" in filtered.columns and filtered["Time Received"].notna().any():
    mn = filtered["Time Received"].min().date()
    mx = filtered["Time Received"].max().date()
    dr = st.sidebar.date_input("Date Range", value=(mn, mx))
    if isinstance(dr, (list, tuple)) and len(dr) == 2:
        filtered = filtered[
            (filtered["Time Received"].dt.date >= dr[0]) &
            (filtered["Time Received"].dt.date <= dr[1])
        ]

st.sidebar.success(f"✅ {len(filtered):,} alarms loaded")

# ── Context bar: active file + quick KPIs ────────────────────────────────────
_n_dev = filtered["Device"].nunique() if "Device" in filtered.columns else 0
_n_biu = filtered["BIU"].nunique() if "BIU" in filtered.columns else 0
if "Time Received" in filtered.columns and filtered["Time Received"].notna().any():
    _span = f"{filtered['Time Received'].min():%b %d} → {filtered['Time Received'].max():%b %d, %Y}"
else:
    _span = "—"
st.markdown(
    f"""
    <div class="context-bar">
      <div class="ctx-chip ctx-active"><div class="cl">Active Source</div>
        <div class="cv">{sel_wb} › {sel_sheet}</div></div>
      <div class="ctx-chip"><div class="cl">Alarms</div><div class="cv">{len(filtered):,}</div></div>
      <div class="ctx-chip"><div class="cl">Devices</div><div class="cv">{_n_dev}</div></div>
      <div class="ctx-chip"><div class="cl">BIUs</div><div class="cv">{_n_biu}</div></div>
      <div class="ctx-chip"><div class="cl">Date Range</div><div class="cv">{_span}</div></div>
    </div>
    """,
    unsafe_allow_html=True,
)


# ════════════════════════════════════════════════════════════════════════════
# TABS
# ════════════════════════════════════════════════════════════════════════════
tab_ov, tab_an, tab_hlt, tab_tm, tab_pred, tab_ml, tab_cmp = st.tabs([
    "📊 Overview", "🔍 Alarm Analysis", "🏥 Device Health",
    "🕐 Time Patterns", "🔮 Predictions", "🤖 ML Forecast", "⚖️ Compare Files",
])


# ════════════════════════════════════════════════════════════════════════════
# TAB 1 — OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
with tab_ov:
    st.subheader(f"📊  {sel_wb}  ›  {sel_sheet}")
    c1,c2,c3,c4,c5 = st.columns(5)
    c1.metric("Total Alarms",   f"{len(filtered):,}")
    c2.metric("Unique Levels",  filtered["Level"].nunique()  if "Level"  in filtered.columns else 0)
    c3.metric("Unique BIUs",    filtered["BIU"].nunique()    if "BIU"    in filtered.columns else 0)
    c4.metric("Unique Devices", filtered["Device"].nunique() if "Device" in filtered.columns else 0)
    ad = filtered["Duration (Min)"].dropna().mean() if "Duration (Min)" in filtered.columns else 0
    c5.metric("Avg Duration (Min)", round(float(ad),1) if not pd.isna(ad) else 0)

    st.markdown("---")
    col1, col2 = st.columns(2)

    if "Level" in filtered.columns and not filtered.empty:
        lc = filtered["Level"].value_counts().reset_index(); lc.columns=["Level","Count"]
        fig = px.bar(lc, x="Level", y="Count", text="Count", title="Alarms by Level",
                     color="Level", color_discrete_map=LEVEL_COLOR)
        fig.update_traces(textposition="outside", texttemplate="%{text:.0f}")
        fig.update_layout(showlegend=False, yaxis=dict(tickformat="d"))
        col1.plotly_chart(apply_theme(fig), use_container_width=True, config=PCFG)

    if "BIU" in filtered.columns and not filtered.empty:
        bc = filtered["BIU"].value_counts().reset_index(); bc.columns=["BIU","Count"]
        fig2 = px.bar(bc, x="BIU", y="Count", text="Count", title="Alarms by BIU")
        fig2.update_traces(textposition="outside", texttemplate="%{text:.0f}")
        fig2.update_layout(yaxis=dict(tickformat="d"))
        col2.plotly_chart(apply_theme(fig2), use_container_width=True, config=PCFG)

    if "Date" in filtered.columns and not filtered.empty:
        tc = filtered.groupby("Date").size().reset_index(name="Count")
        fig3 = px.line(tc, x="Date", y="Count", markers=True, title="Daily Alarm Timeline")
        fig3.update_traces(line_color=blue,
                           hovertemplate="Date: %{x}<br>Alarms: %{y:d}<extra></extra>")
        fig3.update_layout(yaxis=dict(tickformat="d"))
        st.plotly_chart(apply_theme(fig3), use_container_width=True, config=PCFG)
        dl_btn(fig3, "daily_timeline.png", "Download Timeline")

    with st.expander("🗂️  Raw Data Preview"):
        st.dataframe(filtered, use_container_width=True)


# ════════════════════════════════════════════════════════════════════════════
# TAB 2 — ALARM ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
with tab_an:
    if filtered.empty:
        st.warning("No data after filtering.")
    else:
        # Level pie + duration box
        if "Level" in filtered.columns:
            col1, col2 = st.columns(2)
            pied = filtered["Level"].value_counts().reset_index(); pied.columns=["Level","Count"]
            fig_pie = px.pie(pied, names="Level", values="Count", title="Level Distribution",
                             hole=0.35, color="Level", color_discrete_map=LEVEL_COLOR)
            col1.plotly_chart(apply_theme(fig_pie), use_container_width=True, config=PCFG)
            dl_btn(fig_pie, "level_pie.png", "Download Pie Chart")

            if "Duration (Min)" in filtered.columns:
                fig_box = px.box(
                    filtered.dropna(subset=["Duration (Min)"]),
                    x="Level", y="Duration (Min)", title="Alarm Duration by Level",
                    color="Level", color_discrete_map=LEVEL_COLOR)
                fig_box.update_layout(showlegend=False)
                col2.plotly_chart(apply_theme(fig_box), use_container_width=True, config=PCFG)

        # Top devices — short name display, full path in hover
        if "Device" in filtered.columns and not filtered.empty:
            st.subheader("Top Affected Devices")
            st.caption("Hover a bar to see the full device path.")
            dev_g = (
                filtered.groupby(["Device Short","Device"])
                .agg(Count=("Device","count"), Avg_Sev=("Severity","mean"))
                .reset_index().sort_values("Count", ascending=False).head(20)
            )
            fig_dev = px.bar(dev_g, x="Device Short", y="Count", text="Count",
                             title="Alarm Count by Device",
                             custom_data=["Device","Avg_Sev"],
                             color="Avg_Sev",
                             color_continuous_scale=["#43a047","#fdd835","#fb8c00","#e53935"],
                             labels={"Avg_Sev":"Avg Severity"})
            fig_dev.update_traces(
                textposition="outside", texttemplate="%{text:.0f}",
                hovertemplate=(
                    "<b>%{x}</b><br>📍 Full path: %{customdata[0]}<br>"
                    "Count: %{y:d}<br>Avg Severity: %{customdata[1]:.2f}<extra></extra>"
                ))
            fig_dev.update_layout(coloraxis_colorbar=dict(title="Avg Severity"),
                                  yaxis=dict(tickformat="d"))
            st.plotly_chart(apply_theme(fig_dev), use_container_width=True, config=PCFG)
            dl_btn(fig_dev, "top_devices.png", "Download Device Chart")

        # Level × BIU heatmap
        if "Level" in filtered.columns and "BIU" in filtered.columns:
            st.subheader("Level × BIU Heatmap")
            heat = filtered.groupby(["Level","BIU"]).size().reset_index(name="Count")
            hp   = heat.pivot(index="Level", columns="BIU", values="Count").fillna(0).astype(int)
            mxv  = int(hp.values.max()) if hp.size > 0 else 1
            fig_heat = go.Figure(go.Heatmap(
                z=hp.values, x=hp.columns.tolist(), y=hp.index.tolist(),
                text=hp.values, texttemplate="%{text:d}",
                colorscale=[[0,"#fffde7"],[0.33,"#fdd835"],[0.66,"#fb8c00"],[1,"#e53935"]],
                zmin=0, zmax=mxv,
                hovertemplate="Level: %{y}<br>BIU: %{x}<br>Count: %{z:d}<extra></extra>"))
            fig_heat.update_layout(title="Alarm Concentration: Level vs BIU")
            st.plotly_chart(apply_theme(fig_heat), use_container_width=True, config=PCFG)
            dl_btn(fig_heat, "level_biu_heatmap.png", "Download Heatmap")

        # Device × BIU heatmap
        if "Device Short" in filtered.columns and "BIU" in filtered.columns:
            st.subheader("Device × BIU Heatmap")
            dh = filtered.groupby(["Device Short","BIU"]).size().reset_index(name="Count")
            top15 = filtered["Device Short"].value_counts().head(15).index
            dh    = dh[dh["Device Short"].isin(top15)]
            dp    = dh.pivot(index="Device Short", columns="BIU", values="Count").fillna(0).astype(int)
            fig_dh = go.Figure(go.Heatmap(
                z=dp.values, x=dp.columns.tolist(), y=dp.index.tolist(),
                text=dp.values, texttemplate="%{text:d}",
                colorscale=[[0,"#161b22"],[0.5,blue],[1,"#e53935"]],
                hovertemplate="Device: %{y}<br>BIU: %{x}<br>Count: %{z:d}<extra></extra>"))
            fig_dh.update_layout(title="Device vs BIU Alarm Concentration (Top 15 Devices)")
            st.plotly_chart(apply_theme(fig_dh), use_container_width=True, config=PCFG)

        # Top alarm names + messages
        col3, col4 = st.columns(2)
        if "Name" in filtered.columns and not filtered.empty:
            nc = filtered["Name"].value_counts().reset_index().head(15); nc.columns=["Name","Count"]
            fn = px.bar(nc, x="Count", y="Name", orientation="h", text="Count", title="Top Alarm Names")
            fn.update_layout(yaxis={"categoryorder":"total ascending"}, xaxis=dict(tickformat="d"))
            col3.plotly_chart(apply_theme(fn), use_container_width=True, config=PCFG)

        if "Message" in filtered.columns and not filtered.empty:
            mc = filtered["Message"].value_counts().reset_index().head(10); mc.columns=["Message","Count"]
            fm = px.bar(mc, x="Count", y="Message", orientation="h", text="Count", title="Top Messages")
            fm.update_layout(yaxis={"categoryorder":"total ascending"}, xaxis=dict(tickformat="d"))
            col4.plotly_chart(apply_theme(fm), use_container_width=True, config=PCFG)

        # Sunburst hierarchy
        h_cols = [c for c in ["Level","BIU","Device Short","Name"] if c in filtered.columns]
        if len(h_cols) >= 2 and not filtered.empty:
            st.subheader("Hierarchy View  —  Level → BIU → Device → Alarm Name")
            sun = filtered.groupby(h_cols).size().reset_index(name="Count")
            fig_sun = px.sunburst(sun, path=h_cols, values="Count",
                                  title="Level → BIU → Device → Name")
            st.plotly_chart(apply_theme(fig_sun), use_container_width=True, config=PCFG)

        # Duration histogram
        if "Duration (Min)" in filtered.columns and filtered["Duration (Min)"].notna().any():
            st.subheader("Alarm Duration Distribution")
            fig_dur = px.histogram(filtered.dropna(subset=["Duration (Min)"]),
                                   x="Duration (Min)", nbins=30, title="Duration Distribution",
                                   color_discrete_sequence=[blue])
            st.plotly_chart(apply_theme(fig_dur), use_container_width=True, config=PCFG)

        # PPT export
        st.markdown("---")
        st.subheader("📥 Export PowerPoint Report")
        if st.button("Generate PowerPoint Report"):
            try:
                tpl = ppt_tpl.getvalue() if ppt_tpl else None
                buf = build_ppt(filtered, sel_wb, sel_sheet, tpl)
                st.download_button("⬇️  Download PowerPoint", buf,
                    f"Alarm_Report_{sel_wb}_{sel_sheet}.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                st.success("PPT ready. Charts are editable inside PowerPoint.")
            except Exception as e:
                st.error(f"PPT error: {e}")


# ════════════════════════════════════════════════════════════════════════════
# TAB 3 — DEVICE HEALTH
# ════════════════════════════════════════════════════════════════════════════
with tab_hlt:
    st.subheader("🏥 Device Health Scorecard")
    st.caption(
        "Health Score 0–100:  100 = healthy,  0 = critical.  "
        "Weighted: alarm frequency (35%) + severity (35%) + MTTR (15%) + MTBF (15%)."
    )
    hdf = compute_health(filtered)

    if hdf.empty:
        st.warning("Not enough data for health analysis.")
    else:
        c1,c2,c3,c4 = st.columns(4)
        c1.metric("🔴 Critical", len(hdf[hdf["Health Score"] < 40]))
        c2.metric("🟠 Poor",     len(hdf[(hdf["Health Score"] >= 40) & (hdf["Health Score"] < 60)]))
        c3.metric("🟡 Fair",     len(hdf[(hdf["Health Score"] >= 60) & (hdf["Health Score"] < 80)]))
        c4.metric("🟢 Good",     len(hdf[hdf["Health Score"] >= 80]))

        # Health bar chart (top 20 at-risk)
        top20 = hdf.head(20).copy()
        fig_h = px.bar(top20, x="Health Score", y="Device", orientation="h",
                       text="Health Score",
                       title="Device Health Scores — Top 20 Most At-Risk",
                       color="Health Score",
                       color_continuous_scale=["#e53935","#fb8c00","#fdd835","#43a047"],
                       range_color=[0,100],
                       custom_data=["Device Full","Alarms","MTBF (Hrs)","MTTR (Min)","Top Level","BIU"])
        fig_h.update_traces(
            texttemplate="%{text:.1f}", textposition="outside",
            hovertemplate=(
                "<b>%{y}</b><br>📍 Full path: %{customdata[0]}<br>"
                "BIU: %{customdata[5]}<br>Health Score: %{x:.1f}<br>"
                "Total Alarms: %{customdata[1]}<br>MTBF: %{customdata[2]} hrs<br>"
                "MTTR: %{customdata[3]} min<br>Top Level: %{customdata[4]}<extra></extra>"
            ))
        fig_h.update_layout(yaxis={"categoryorder":"total ascending"},
                            coloraxis_colorbar=dict(title="Score"),
                            xaxis=dict(range=[0,115]))
        st.plotly_chart(apply_theme(fig_h), use_container_width=True, config=PCFG)
        dl_btn(fig_h, "device_health.png", "Download Health Chart")

        # MTBF vs MTTR reliability scatter
        if hdf["MTBF (Hrs)"].notna().any():
            st.subheader("MTBF vs MTTR — Reliability Map")
            st.caption("Best: far-right (long between failures) + low (fast recovery).  Worst: bottom-left.")
            fig_sc = px.scatter(
                hdf.dropna(subset=["MTBF (Hrs)"]),
                x="MTBF (Hrs)", y="MTTR (Min)",
                size="Alarms", color="Health Score",
                color_continuous_scale=["#e53935","#fb8c00","#fdd835","#43a047"],
                range_color=[0,100], text="Device",
                title="Reliability Map: MTBF vs MTTR",
                custom_data=["Device Full","Alarms","Top Level"])
            fig_sc.update_traces(
                textposition="top center",
                hovertemplate=(
                    "<b>%{text}</b><br>📍 %{customdata[0]}<br>"
                    "MTBF: %{x} hrs  |  MTTR: %{y} min<br>"
                    "Alarms: %{customdata[1]}<br>Top Level: %{customdata[2]}<extra></extra>"
                ))
            fig_sc.update_layout(
                xaxis_title="MTBF (hours)  ←  frequent failures          rare failures  →",
                yaxis_title="MTTR (minutes)  ←  fast recovery          slow recovery  →")
            st.plotly_chart(apply_theme(fig_sc), use_container_width=True, config=PCFG)

        # Full health table
        st.subheader("Complete Device Health Table")
        show = ["Device","BIU","Status","Health Score","Alarms","Alarms/Day",
                "Top Level","MTBF (Hrs)","MTTR (Min)","Last Alarm","Top Alarm"]
        disp = hdf[[c for c in show if c in hdf.columns]].copy()
        if "Last Alarm" in disp.columns:
            disp["Last Alarm"] = disp["Last Alarm"].dt.strftime("%Y-%m-%d %H:%M")
        st.dataframe(disp, use_container_width=True)
        st.download_button("⬇️  Download Health Table (CSV)",
                           disp.to_csv(index=False).encode(),
                           "device_health.csv", "text/csv")


# ════════════════════════════════════════════════════════════════════════════
# TAB 4 — TIME PATTERNS
# ════════════════════════════════════════════════════════════════════════════
with tab_tm:
    st.subheader("🕐 Temporal Alarm Patterns")
    if "Hour" not in filtered.columns or filtered.empty:
        st.warning("No timestamp data available.")
    else:
        DOW = ["Monday","Tuesday","Wednesday","Thursday","Friday","Saturday","Sunday"]
        col1, col2 = st.columns(2)

        # Alarms by hour
        hc = filtered.groupby("Hour").size().reset_index(name="Count")
        fig_hr = px.bar(hc, x="Hour", y="Count", text="Count",
                        title="Alarms by Hour of Day", color="Count",
                        color_continuous_scale=[blue,"#e53935"])
        fig_hr.update_traces(textposition="outside", texttemplate="%{text:.0f}")
        fig_hr.update_layout(xaxis=dict(tickmode="linear",tick0=0,dtick=1),
                             yaxis=dict(tickformat="d"))
        col1.plotly_chart(apply_theme(fig_hr), use_container_width=True, config=PCFG)

        # Alarms by day of week
        dw = filtered.groupby("DayOfWeek").size().reset_index(name="Count")
        dw["DayOfWeek"] = pd.Categorical(dw["DayOfWeek"], categories=DOW, ordered=True)
        dw = dw.sort_values("DayOfWeek")
        fig_dw = px.bar(dw, x="DayOfWeek", y="Count", text="Count",
                        title="Alarms by Day of Week", color="Count",
                        color_continuous_scale=[blue,"#e53935"])
        fig_dw.update_traces(textposition="outside", texttemplate="%{text:.0f}")
        fig_dw.update_layout(yaxis=dict(tickformat="d"))
        col2.plotly_chart(apply_theme(fig_dw), use_container_width=True, config=PCFG)

        # Hour × Day heatmap
        st.subheader("Peak Hours Heatmap")
        hd = filtered.groupby(["DayOfWeek","Hour"]).size().reset_index(name="Count")
        hd["DayOfWeek"] = pd.Categorical(hd["DayOfWeek"], categories=DOW, ordered=True)
        hp = hd.pivot(index="DayOfWeek", columns="Hour", values="Count").fillna(0).astype(int)
        fig_hm = go.Figure(go.Heatmap(
            z=hp.values, x=list(hp.columns), y=list(hp.index),
            text=hp.values, texttemplate="%{text:d}",
            colorscale=[[0,"#161b22"],[0.5,blue],[1,"#e53935"]],
            hovertemplate="Day: %{y}<br>Hour: %{x}:00<br>Alarms: %{z:d}<extra></extra>",
            colorbar=dict(title="Alarms")))
        fig_hm.update_layout(title="When Do Alarms Peak? (Day × Hour)",
                             xaxis_title="Hour of Day (0 = midnight)")
        st.plotly_chart(apply_theme(fig_hm), use_container_width=True, config=PCFG)
        dl_btn(fig_hm, "peak_heatmap.png", "Download Peak Heatmap")

        # Severity by hour stacked
        if "Level" in filtered.columns:
            st.subheader("Severity Distribution by Hour")
            hl = filtered.groupby(["Hour","Level"]).size().reset_index(name="Count")
            fig_hl = px.bar(hl, x="Hour", y="Count", color="Level",
                            title="Alarm Severity by Hour",
                            color_discrete_map=LEVEL_COLOR, barmode="stack")
            fig_hl.update_layout(xaxis=dict(tickmode="linear",tick0=0,dtick=1),
                                 yaxis=dict(tickformat="d"))
            st.plotly_chart(apply_theme(fig_hl), use_container_width=True, config=PCFG)

        # BIU trend over time
        if "BIU" in filtered.columns and "Date" in filtered.columns:
            st.subheader("BIU Alarm Trend Over Time")
            bt = filtered.groupby(["Date","BIU"]).size().reset_index(name="Count")
            fig_bt = px.line(bt, x="Date", y="Count", color="BIU",
                             title="Daily Alarms per BIU", markers=True)
            fig_bt.update_layout(yaxis=dict(tickformat="d"))
            st.plotly_chart(apply_theme(fig_bt), use_container_width=True, config=PCFG)


# ════════════════════════════════════════════════════════════════════════════
# TAB 5 — PREDICTIONS
# ════════════════════════════════════════════════════════════════════════════
with tab_pred:
    st.subheader("🔮 Alarm Prediction Engine")
    st.caption(
        "Predictions use exponentially-weighted mean inter-arrival time per device.  "
        "Higher Confidence % = more regular, predictable alarm pattern."
    )

    with st.expander("ℹ️  How predictions work (no ML training required)"):
        st.markdown(
            "This engine is a **live statistical estimate**, not a trained model — it is "
            "recomputed from the uploaded data every time, so there is nothing to train.\n\n"
            "For each device:\n"
            "1. **Inter-arrival times** — the gaps (in hours) between consecutive alarms.\n"
            "2. **Exponentially-weighted mean (EWM)** of those gaps — recent gaps count more, "
            "so the estimate adapts to a changing alarm rate.\n"
            "3. **Predicted next alarm** = *last alarm time* + *EWM gap*.\n"
            "4. **Confidence %** = `100 × (1 − CV)`, where CV is the coefficient of variation "
            "(std ÷ mean) of the gaps. A regular, evenly-spaced pattern → high confidence; "
            "an erratic one → low confidence.\n\n"
            "👉 For a *trained* model that forecasts future alarm volume, see the "
            "**🤖 ML Forecast** tab."
        )

    hdf_p = compute_health(filtered)

    # ── Reference time ────────────────────────────────────────────────────────
    # Urgency must be measured against a sensible "now". When the uploaded data is
    # historical (e.g. months old), datetime.now() makes EVERY device look overdue,
    # so default the anchor to the latest alarm actually present in the data.
    data_latest = (
        filtered["Time Received"].max()
        if "Time Received" in filtered.columns and filtered["Time Received"].notna().any()
        else None
    )
    data_earliest = (
        filtered["Time Received"].min()
        if "Time Received" in filtered.columns and filtered["Time Received"].notna().any()
        else None
    )

    anchor_choice = st.radio(
        "Measure predictions relative to:",
        ["Latest alarm in data", "Now (wall clock)"],
        index=0, horizontal=True,
        help="Historical exports are usually best measured from the last alarm in the file.",
    )
    if anchor_choice == "Now (wall clock)" or data_latest is None:
        ref_time = datetime.now()
    else:
        ref_time = data_latest.to_pydatetime() if hasattr(data_latest, "to_pydatetime") else data_latest

    if data_earliest is not None and data_latest is not None:
        st.caption(
            f"📅 Data spans **{data_earliest:%Y-%m-%d %H:%M}** → **{data_latest:%Y-%m-%d %H:%M}**  ·  "
            f"reference time = **{ref_time:%Y-%m-%d %H:%M}** ({anchor_choice})."
        )

    if hdf_p.empty or "Predicted Next" not in hdf_p.columns:
        st.warning("Need at least 2 recorded alarms per device to generate predictions.")
    else:
        pred = hdf_p.dropna(subset=["Predicted Next"]).copy()
        pred["Hours Until Next"] = (
            (pred["Predicted Next"] - ref_time).dt.total_seconds() / 3600
        ).round(1)

        def urgency(h):
            if h < 0:    return "⚠️ Overdue"
            if h < 24:   return "🔴 Imminent (<24h)"
            if h < 48:   return "🟠 Soon (24–48h)"
            if h < 72:   return "🟡 Upcoming (48–72h)"
            return "🟢 Later (>72h)"

        pred["Urgency"] = pred["Hours Until Next"].apply(urgency)

        urg = pred["Urgency"].value_counts()
        c1,c2,c3,c4,c5 = st.columns(5)
        c1.metric("⚠️ Overdue",          urg.get("⚠️ Overdue",           0))
        c2.metric("🔴 Imminent (<24h)",  urg.get("🔴 Imminent (<24h)",   0))
        c3.metric("🟠 Soon (24–48h)",    urg.get("🟠 Soon (24–48h)",     0))
        c4.metric("🟡 Upcoming (48–72h)",urg.get("🟡 Upcoming (48–72h)", 0))
        c5.metric("🟢 Later (>72h)",     urg.get("🟢 Later (>72h)",      0))

        # Prediction timeline
        st.subheader("Predicted Alarm Timeline")
        COLOR_URG = {
            "⚠️ Overdue":           "#9c27b0",
            "🔴 Imminent (<24h)":   "#e53935",
            "🟠 Soon (24–48h)":     "#fb8c00",
            "🟡 Upcoming (48–72h)": "#fdd835",
            "🟢 Later (>72h)":      "#43a047",
        }
        fig_pred = px.scatter(
            pred.sort_values("Predicted Next").head(40),
            x="Predicted Next", y="Device",
            color="Urgency", size="Confidence (%)",
            color_discrete_map=COLOR_URG,
            title="Next Predicted Alarm per Device",
            custom_data=["Device Full","MTBF (Hrs)","Confidence (%)","Top Level","Hours Until Next"])
        fig_pred.update_traces(
            hovertemplate=(
                "<b>%{y}</b><br>📍 %{customdata[0]}<br>"
                "Predicted: %{x}<br>Hours until: %{customdata[4]:.1f} h<br>"
                "MTBF: %{customdata[1]} hrs<br>Confidence: %{customdata[2]}%%<br>"
                "Typical Level: %{customdata[3]}<extra></extra>"
            ))
        # NOTE: add_vline with a string x + annotation_text crashes on plotly 6.x
        # (it averages the x-coords → int + str TypeError). Draw the line with a
        # real datetime and add the label as a separate annotation instead.
        fig_pred.add_vline(x=ref_time, line_dash="dash", line_color="gray")
        fig_pred.add_annotation(
            x=ref_time, yref="paper", y=1.0, yanchor="bottom",
            text="⏱ Reference", showarrow=False,
            font=dict(color="gray", size=12),
        )
        fig_pred.update_layout(yaxis={"categoryorder":"total ascending"})
        st.plotly_chart(apply_theme(fig_pred), use_container_width=True, config=PCFG)

        # Prediction table
        st.subheader("Full Prediction Table")
        ptbl = pred[["Device","BIU","Urgency","Predicted Next","Hours Until Next",
                     "Confidence (%)","MTBF (Hrs)","Top Level","Alarms"]].copy()
        ptbl["Predicted Next"] = ptbl["Predicted Next"].dt.strftime("%Y-%m-%d %H:%M")
        ptbl = ptbl.sort_values("Hours Until Next")
        st.dataframe(ptbl, use_container_width=True)
        st.download_button("⬇️  Download Predictions (CSV)",
                           ptbl.to_csv(index=False).encode(),
                           "alarm_predictions.csv", "text/csv")

        # Device risk score bar
        st.subheader("Device Risk Ranking")
        risk = pred.copy()
        risk["Risk Score"] = (
            risk["Alarms/Day"] * risk["Avg Severity"] * (101 - risk["Health Score"])
        ).round(2)
        risk = risk.sort_values("Risk Score", ascending=False).head(15)
        fig_risk = px.bar(risk, x="Risk Score", y="Device", orientation="h",
                          text="Risk Score", title="Device Risk Score (frequency × severity × health penalty)",
                          color="Risk Score",
                          color_continuous_scale=["#43a047","#fdd835","#fb8c00","#e53935"],
                          custom_data=["Device Full","Health Score","Top Level"])
        fig_risk.update_traces(
            texttemplate="%{text:.1f}", textposition="outside",
            hovertemplate=(
                "<b>%{y}</b><br>📍 %{customdata[0]}<br>"
                "Risk Score: %{x:.1f}<br>Health: %{customdata[1]}<br>"
                "Top Level: %{customdata[2]}<extra></extra>"
            ))
        fig_risk.update_layout(yaxis={"categoryorder":"total ascending"})
        st.plotly_chart(apply_theme(fig_risk), use_container_width=True, config=PCFG)


# ════════════════════════════════════════════════════════════════════════════
# TAB 6 — ML FORECAST  (trained model)
# ════════════════════════════════════════════════════════════════════════════
with tab_ml:
    st.markdown('<div class="section-title">🤖 ML Alarm-Volume Forecast</div>',
                unsafe_allow_html=True)
    st.caption(
        "A **trained** Gradient Boosting model that learns daily alarm patterns "
        "(day-of-week, recent lags &amp; rolling averages) and forecasts upcoming alarm volume. "
        "Unlike the 🔮 Predictions tab (a live statistic), this model is fit on your data and can be exported."
    )

    if not SKLEARN_AVAILABLE:
        st.error("scikit-learn is required for ML forecasting.  Install it with:  "
                 "`pip install scikit-learn`")
    else:
        # ── Training data: current sheet + any extra historical files ─────────
        st.markdown('<div class="section-title">📈 Training Data</div>',
                    unsafe_allow_html=True)
        st.caption(
            "The model trains on the **current sheet's full history** (sidebar filters are "
            "ignored here). Add more **historical EventLog files for the *same site*** below to "
            "extend the training window — more continuous history → better accuracy."
        )
        extra_files = st.file_uploader(
            "➕ Add historical alarm files to train on (.xls / .xlsx)",
            type=["xls", "xlsx"], accept_multiple_files=True, key="ml_train_files",
        )

        frames = [df]                                   # current sheet, unfiltered
        sources = [(f"{sel_wb} › {sel_sheet} (current)", len(df))]
        if extra_files:
            for uf in extra_files:
                try:
                    h = load_alarm_history(uf.getvalue(), uf.name)
                    if h.empty or "Time Received" not in h.columns:
                        st.warning(f"⚠️ No alarm data found in **{uf.name}** — skipped.")
                        continue
                    frames.append(h)
                    sources.append((uf.name, len(h)))
                except Exception as e:
                    st.error(f"Could not read {uf.name}: {e}")

        train_df = combine_histories(frames)
        daily = build_daily_series(train_df)

        # Show what the model is training on
        src_lines = "  ·  ".join(f"{name} ({n:,})" for name, n in sources)
        st.markdown(
            f"**Training on {len(sources)} source(s):** {src_lines}  "
            f"→ **{len(train_df):,}** unique alarms after de-duplication."
        )
        if len(sources) > 1:
            st.info(
                "ℹ️ Combining multiple files trains **one** volume model across all of them. "
                "For a coherent per-site forecast, only add files from the **same** monitored system."
            )
        # Let the user keep a growing master history to re-upload next time
        if "Time Received" in train_df.columns:
            st.download_button(
                "⬇️  Download combined history (CSV — keep & re-upload to accumulate)",
                train_df.to_csv(index=False).encode(),
                f"alarm_training_history_{sel_wb}.csv", "text/csv",
                key="dl_train_hist",
            )

        if daily.empty or len(daily) < 21:
            st.warning(
                f"Need at least **21 days** of alarm history to train a reliable model "
                f"(current training set spans **{len(daily)}** day(s)). "
                "Add more historical files above, or upload a file with a longer date range."
            )
        else:
            st.markdown(
                f"**Training window:** {daily['Date'].min():%Y-%m-%d} → "
                f"{daily['Date'].max():%Y-%m-%d}  ·  {len(daily)} days  ·  "
                f"{int(daily['Count'].sum()):,} total alarms"
            )
            result = train_volume_model(daily)

            if not result.get("ok"):
                st.warning(f"Could not train a model: {result.get('reason')}.")
            else:
                m = result["metrics"]
                c1, c2, c3, c4 = st.columns(4)
                c1.metric("Train / Test days", f"{result['n_train']} / {result['n_test']}")
                c2.metric("MAE", f"{m['MAE']:.2f}", help="Mean absolute error (alarms/day)")
                c3.metric("RMSE", f"{m['RMSE']:.2f}")
                c4.metric("R²", "n/a" if np.isnan(m["R2"]) else f"{m['R2']:.2f}",
                          help="1.0 = perfect; ≤0 = worse than predicting the average")

                # Actual vs predicted on the holdout
                st.markdown('<div class="section-title">Model Accuracy (holdout)</div>',
                            unsafe_allow_html=True)
                ev = result["eval_df"]
                fig_ev = go.Figure()
                fig_ev.add_trace(go.Scatter(x=ev["Date"], y=ev["Actual"], mode="lines+markers",
                                            name="Actual", line=dict(color=blue)))
                fig_ev.add_trace(go.Scatter(x=ev["Date"], y=ev["Predicted"], mode="lines+markers",
                                            name="Predicted", line=dict(color="#e53935", dash="dash")))
                fig_ev.update_layout(title="Actual vs Predicted Daily Alarms (test period)",
                                     yaxis=dict(tickformat="d"))
                st.plotly_chart(apply_theme(fig_ev), use_container_width=True, config=PCFG)

                # Forecast
                st.markdown('<div class="section-title">Forecast</div>', unsafe_allow_html=True)
                horizon = st.slider("Forecast horizon (days)", 7, 30, 14)
                fc = forecast_volume(result["model"], daily, horizon)
                fig_fc = go.Figure()
                fig_fc.add_trace(go.Scatter(x=daily["Date"], y=daily["Count"], mode="lines",
                                            name="History", line=dict(color=blue)))
                fig_fc.add_trace(go.Scatter(x=fc["Date"], y=fc["Forecast"], mode="lines+markers",
                                            name="Forecast", line=dict(color="#fb8c00", dash="dot")))
                fig_fc.add_vline(x=daily["Date"].max(), line_dash="dash", line_color="gray")
                fig_fc.add_annotation(x=daily["Date"].max(), yref="paper", y=1.0, yanchor="bottom",
                                      text="Forecast start", showarrow=False,
                                      font=dict(color="gray", size=12))
                fig_fc.update_layout(title=f"Next {horizon}-Day Alarm-Volume Forecast",
                                     yaxis=dict(tickformat="d"))
                st.plotly_chart(apply_theme(fig_fc), use_container_width=True, config=PCFG)
                st.metric("Forecast total (next %d days)" % horizon, f"{fc['Forecast'].sum():.0f}")

                col_a, col_b = st.columns(2)
                col_a.download_button(
                    "⬇️  Download Forecast (CSV)",
                    fc.assign(Date=fc["Date"].dt.strftime("%Y-%m-%d")).to_csv(index=False).encode(),
                    "alarm_volume_forecast.csv", "text/csv")

                # Export the trained model
                try:
                    mbuf = io.BytesIO()
                    joblib.dump({"model": result["model"], "features": ML_FEATURES}, mbuf)
                    mbuf.seek(0)
                    col_b.download_button(
                        "⬇️  Download Trained Model (.joblib)", mbuf.getvalue(),
                        f"alarm_volume_model_{sel_wb}.joblib", "application/octet-stream")
                except Exception as e:
                    col_b.info(f"Model export unavailable: {e}")

                # Feature importance
                with st.expander("🔬 Feature importance & caveats"):
                    imp = (pd.Series(result["importances"]).sort_values(ascending=True))
                    fig_imp = px.bar(imp, orientation="h", title="What drives the forecast",
                                     labels={"value": "Importance", "index": "Feature"})
                    fig_imp.update_layout(showlegend=False)
                    st.plotly_chart(apply_theme(fig_imp), use_container_width=True, config=PCFG)
                    st.caption(
                        "Caveats: alarm data is often sparse and bursty, so daily-volume forecasts "
                        "are indicative, not exact. The model is trained only on the current selection "
                        "and is re-trained whenever the data or filters change."
                    )


# ════════════════════════════════════════════════════════════════════════════
# TAB 7 — COMPARE FILES
# ════════════════════════════════════════════════════════════════════════════
with tab_cmp:
    st.subheader("⚖️  Multi-File Comparison")
    if len(wmap) < 2:
        st.info("📂  Upload at least 2 Excel files to use comparison mode.")
    else:
        st.markdown("**Select which files and sheets to compare:**")
        selections = []
        for wb_name, wb_data in wmap.items():
            ca, cb = st.columns([2,2])
            with ca:
                inc = st.checkbox(f"Include: {wb_name}", value=True, key=f"inc_{wb_name}")
            with cb:
                sh  = st.selectbox("Sheet", wb_data["sheets"], key=f"sh_{wb_name}")
            if inc:
                selections.append((wb_name, sh, wb_data["file"]))

        if len(selections) < 2:
            st.warning("Select at least 2 files.")
        else:
            comp = {}
            for wb_name, sh, uf in selections:
                label = f"{wb_name} › {sh}"
                try:
                    comp[label] = clean_df(read_sheet(uf.getvalue(), uf.name, sh))
                except Exception as e:
                    st.error(f"Could not load {label}: {e}")

            if len(comp) >= 2:
                # Summary metrics
                st.subheader("Summary Comparison")
                summ = pd.DataFrame({
                    "File":           list(comp.keys()),
                    "Total Alarms":   [len(d) for d in comp.values()],
                    "Avg Duration":   [
                        round(float(d["Duration (Min)"].dropna().mean()), 1)
                        if "Duration (Min)" in d.columns and d["Duration (Min)"].notna().any() else 0
                        for d in comp.values()],
                    "Unique Devices": [d["Device"].nunique() if "Device" in d.columns else 0 for d in comp.values()],
                    "Unique BIUs":    [d["BIU"].nunique()    if "BIU"    in d.columns else 0 for d in comp.values()],
                })
                c1, c2 = st.columns(2)
                fig_ta = px.bar(summ, x="File", y="Total Alarms", text="Total Alarms",
                                title="Total Alarms per File", color="File")
                fig_ta.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                c1.plotly_chart(apply_theme(fig_ta), use_container_width=True, config=PCFG)
                fig_ad = px.bar(summ, x="File", y="Avg Duration", text="Avg Duration",
                                title="Avg Duration (Min) per File", color="File")
                fig_ad.update_traces(textposition="outside", texttemplate="%{text:.1f}")
                c2.plotly_chart(apply_theme(fig_ad), use_container_width=True, config=PCFG)
                st.dataframe(summ, use_container_width=True)

                # Level distribution comparison
                st.subheader("Level Distribution Comparison")
                lc_rows = []
                for label, d in comp.items():
                    if "Level" in d.columns:
                        lc = d["Level"].value_counts().reset_index(); lc.columns=["Level","Count"]; lc["File"]=label
                        lc_rows.append(lc)
                if lc_rows:
                    lc_all = pd.concat(lc_rows)
                    fig_lc = px.bar(lc_all, x="Level", y="Count", color="File",
                                    barmode="group", title="Level Distribution Across Files", text="Count")
                    fig_lc.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                    st.plotly_chart(apply_theme(fig_lc), use_container_width=True, config=PCFG)

                # Top devices comparison
                st.subheader("Top Devices Comparison")
                dc_rows = []
                for label, d in comp.items():
                    if "Device Short" in d.columns:
                        dc = d["Device Short"].value_counts().reset_index().head(10)
                        dc.columns=["Device","Count"]; dc["File"]=label; dc_rows.append(dc)
                if dc_rows:
                    dc_all = pd.concat(dc_rows)
                    fig_dc = px.bar(dc_all, x="Device", y="Count", color="File",
                                    barmode="group", title="Top Devices Across Files", text="Count")
                    fig_dc.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                    st.plotly_chart(apply_theme(fig_dc), use_container_width=True, config=PCFG)

                # Health score comparison
                st.subheader("Device Health Score Comparison")
                hc_rows = []
                for label, d in comp.items():
                    hd = compute_health(d)
                    if not hd.empty:
                        hd["File"] = label
                        hc_rows.append(hd[["Device","Health Score","File","Alarms"]].head(10))
                if hc_rows:
                    hc_all = pd.concat(hc_rows)
                    fig_hc = px.bar(hc_all, x="Device", y="Health Score", color="File",
                                    barmode="group", text="Health Score",
                                    title="Health Scores — Top 10 Most At-Risk Devices",
                                    range_y=[0,115])
                    fig_hc.update_traces(texttemplate="%{text:.1f}", textposition="outside")
                    st.plotly_chart(apply_theme(fig_hc), use_container_width=True, config=PCFG)

                # Daily timeline overlay
                st.subheader("Daily Alarm Timeline Overlay")
                tl_rows = []
                for label, d in comp.items():
                    if "Date" in d.columns:
                        tc = d.groupby("Date").size().reset_index(name="Count"); tc["File"]=label
                        tl_rows.append(tc)
                if tl_rows:
                    tl_all = pd.concat(tl_rows)
                    fig_tl = px.line(tl_all, x="Date", y="Count", color="File",
                                     markers=True, title="Daily Alarms — All Files Overlaid")
                    fig_tl.update_layout(yaxis=dict(tickformat="d"))
                    st.plotly_chart(apply_theme(fig_tl), use_container_width=True, config=PCFG)

                # BIU comparison
                st.subheader("BIU Alarm Comparison")
                biu_rows = []
                for label, d in comp.items():
                    if "BIU" in d.columns:
                        bc = d["BIU"].value_counts().reset_index().head(10)
                        bc.columns=["BIU","Count"]; bc["File"]=label; biu_rows.append(bc)
                if biu_rows:
                    biu_all = pd.concat(biu_rows)
                    fig_biu = px.bar(biu_all, x="BIU", y="Count", color="File",
                                     barmode="group", title="BIU Alarms Across Files", text="Count")
                    fig_biu.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                    st.plotly_chart(apply_theme(fig_biu), use_container_width=True, config=PCFG)

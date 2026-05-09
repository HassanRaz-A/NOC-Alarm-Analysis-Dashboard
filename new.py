import io
from datetime import datetime

import pandas as pd
import streamlit as st
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
import xlrd

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN


# ----------------------------
# Page setup
# ----------------------------
st.set_page_config(page_title="Alarm Dashboard", layout="wide")


# ----------------------------
# Theme selector
# ----------------------------
if "theme_mode" not in st.session_state:
    st.session_state.theme_mode = "Dark"

theme_mode = st.sidebar.selectbox(
    "Select Theme",
    ["Dark", "Light"],
    index=0 if st.session_state.theme_mode == "Dark" else 1
)
st.session_state.theme_mode = theme_mode

if theme_mode == "Dark":
    bg_color = "#0d1117"
    card_color = "#161b22"
    border_color = "#30363d"
    text_color = "#c9d1d9"
    muted_text = "#8b949e"
    neon_blue = "#1e90ff"
    neon_blue_soft = "rgba(30, 144, 255, 0.35)"
    neon_purple = "#8900ff"
    neon_purple_soft = "rgba(137, 0, 255, 0.22)"
    neon_green = "#00ff85"
    plot_template = "plotly_dark"
else:
    bg_color = "#f8fafc"
    card_color = "#ffffff"
    border_color = "#d0d7de"
    text_color = "#111827"
    muted_text = "#4b5563"
    neon_blue = "#2563eb"
    neon_blue_soft = "rgba(37, 99, 235, 0.18)"
    neon_purple = "#7c3aed"
    neon_purple_soft = "rgba(124, 58, 237, 0.14)"
    neon_green = "#10b981"
    plot_template = "plotly_white"


# ----------------------------
# Custom CSS - GitHub + Gemini style
# ----------------------------
st.markdown(
    f"""
    <style>
    :root {{
      --bg-main: {bg_color};
      --bg-card: {card_color};
      --border-subtle: {border_color};
      --text-primary: {text_color};
      --text-muted: {muted_text};
      --neon-blue: {neon_blue};
      --neon-purple: {neon_purple};
      --neon-green: {neon_green};
      --gemini-glow: {neon_blue_soft};
      --gemini-glow-purple: {neon_purple_soft};
    }}

    .stApp {{
        background:
            radial-gradient(circle at top right, var(--gemini-glow-purple), transparent 22%),
            radial-gradient(circle at top left, var(--gemini-glow), transparent 20%),
            var(--bg-main);
        color: var(--text-primary);
    }}

    html, body, [class*="css"] {{
        font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Helvetica, Arial, sans-serif;
    }}

    h1, h2, h3, h4, h5, h6, p, div, span, label {{
        color: var(--text-primary) !important;
    }}

    .block-container {{
        padding-top: 1.1rem;
        padding-bottom: 1.2rem;
    }}

    [data-testid="stSidebar"] {{
        background:
            linear-gradient(180deg, rgba(30,144,255,0.05), rgba(137,0,255,0.03)),
            var(--bg-card);
        border-right: 1px solid var(--border-subtle);
    }}

    /* Cards / metrics */
    div[data-testid="stMetric"] {{
        background:
            linear-gradient(180deg, rgba(30,144,255,0.04), rgba(137,0,255,0.02)),
            var(--bg-card);
        border: 1px solid var(--border-subtle);
        border-radius: 14px;
        padding: 14px;
        transition: all 0.28s ease;
        box-shadow: 0 0 0 rgba(0,0,0,0);
    }}

    div[data-testid="stMetric"]:hover {{
        border-color: var(--neon-blue);
        box-shadow:
            0 0 18px var(--gemini-glow),
            0 0 6px rgba(30, 144, 255, 0.18);
        transform: translateY(-2px);
    }}

    /* Dataframe / code-like panels */
    div[data-testid="stDataFrame"],
    div[data-testid="stTable"] {{
        background-color: var(--bg-card);
        border: 1px solid var(--border-subtle);
        border-radius: 14px;
        padding: 6px;
        transition: all 0.28s ease;
    }}

    div[data-testid="stDataFrame"]:hover,
    div[data-testid="stTable"]:hover {{
        border-color: var(--neon-blue);
        box-shadow:
            0 0 16px var(--gemini-glow),
            0 0 5px rgba(30, 144, 255, 0.16);
    }}

    /* Buttons */
    .stButton > button,
    .stDownloadButton > button {{
        background: transparent !important;
        color: var(--neon-blue) !important;
        border: 2px solid var(--neon-blue) !important;
        border-radius: 10px !important;
        font-weight: 700 !important;
        transition: all 0.22s ease-in-out !important;
        box-shadow: none !important;
    }}

    .stButton > button:hover,
    .stDownloadButton > button:hover {{
        background: var(--neon-blue) !important;
        color: white !important;
        border-color: var(--neon-blue) !important;
        box-shadow: 0 0 20px var(--neon-blue) !important;
        transform: translateY(-1px);
    }}

    /* Inputs */
    .stTextInput > div > div > input,
    .stNumberInput input,
    .stDateInput input,
    .stTimeInput input,
    .stTextArea textarea {{
        background-color: var(--bg-card) !important;
        color: var(--text-primary) !important;
        border: 1px solid var(--border-subtle) !important;
        border-radius: 10px !important;
    }}

    /* Selectbox / multiselect */
    div[data-baseweb="select"] > div {{
        background-color: var(--bg-card) !important;
        border: 1px solid var(--border-subtle) !important;
        border-radius: 10px !important;
        transition: all 0.28s ease !important;
    }}

    div[data-baseweb="select"] > div:hover {{
        border-color: var(--neon-blue) !important;
        box-shadow:
            0 0 14px var(--gemini-glow),
            0 0 5px rgba(30, 144, 255, 0.14) !important;
    }}

    /* File uploader */
    [data-testid="stFileUploader"] {{
        background:
            linear-gradient(180deg, rgba(30,144,255,0.03), rgba(137,0,255,0.02)),
            var(--bg-card);
        border: 1px solid var(--border-subtle);
        border-radius: 14px;
        padding: 10px;
        transition: all 0.28s ease;
    }}

    [data-testid="stFileUploader"]:hover {{
        border-color: var(--neon-blue);
        box-shadow:
            0 0 18px var(--gemini-glow),
            0 0 6px rgba(30, 144, 255, 0.18);
    }}

    /* Expanders */
    details {{
        background-color: var(--bg-card);
        border: 1px solid var(--border-subtle);
        border-radius: 12px;
        padding: 0.35rem 0.75rem;
    }}

    details:hover {{
        border-color: var(--neon-blue);
        box-shadow: 0 0 14px var(--gemini-glow);
    }}

    /* Horizontal rule */
    hr {{
        border-color: var(--border-subtle);
    }}

    /* Section label glow */
    h1 {{
        text-shadow: 0 0 18px {neon_blue_soft};
    }}

    h2, h3 {{
        text-shadow: 0 0 10px rgba(30, 144, 255, 0.08);
    }}

    /* Tabs / pills if any */
    button[kind="secondary"] {{
        border-radius: 10px !important;
    }}
    </style>
    """,
    unsafe_allow_html=True
)

st.title("Alarm Data Visualization Dashboard")
st.markdown("Upload multiple Excel files, choose a sheet, visualize charts, and download an editable PowerPoint report.")


# ----------------------------
# Helpers
# ----------------------------
def apply_theme(fig):
    fig.update_layout(
        template=plot_template,
        paper_bgcolor=bg_color,
        plot_bgcolor=card_color if theme_mode == "Dark" else bg_color,
        font=dict(color=text_color),
        title_font=dict(color=text_color, size=20),
        legend=dict(font=dict(color=text_color)),
        margin=dict(l=40, r=30, t=60, b=40)
    )
    return fig


def load_workbook_sheets(uploaded_file):
    file_name = uploaded_file.name.lower()
    file_bytes = uploaded_file.getvalue()

    if file_name.endswith(".xlsx"):
        xl = pd.ExcelFile(io.BytesIO(file_bytes), engine="openpyxl")
        return xl.sheet_names

    if file_name.endswith(".xls"):
        try:
            xl = pd.ExcelFile(io.BytesIO(file_bytes), engine="xlrd")
            return xl.sheet_names
        except Exception:
            book = xlrd.open_workbook(
                file_contents=file_bytes,
                ignore_workbook_corruption=True
            )
            return book.sheet_names()

    raise ValueError("Unsupported file format. Please upload .xls or .xlsx file.")


def load_selected_sheet(uploaded_file, sheet_name):
    file_name = uploaded_file.name.lower()
    file_bytes = uploaded_file.getvalue()

    try:
        return pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet_name, engine="calamine")
    except Exception:
        pass

    if file_name.endswith(".xlsx"):
        return pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet_name, engine="openpyxl")

    if file_name.endswith(".xls"):
        try:
            return pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet_name, engine="xlrd")
        except Exception:
            book = xlrd.open_workbook(
                file_contents=file_bytes,
                ignore_workbook_corruption=True
            )
            sheet = book.sheet_by_name(sheet_name)
            data = [sheet.row_values(r) for r in range(sheet.nrows)]
            if not data:
                raise ValueError("Selected sheet is empty.")
            headers = [str(h).strip() for h in data[0]]
            return pd.DataFrame(data[1:], columns=headers)

    raise ValueError("Unsupported file format. Please upload .xls or .xlsx file.")


def clean_dataframe(df):
    df = df.copy()
    df.columns = df.columns.astype(str).str.strip()

    for col in ["Type", "Note"]:
        if col in df.columns:
            df.drop(columns=col, inplace=True)

    if "Time Received" in df.columns:
        df["Time Received"] = pd.to_datetime(df["Time Received"], errors="coerce")

    if "Time Cleared" in df.columns:
        df["Time Cleared"] = pd.to_datetime(df["Time Cleared"], errors="coerce")

    if "Time Received" in df.columns and "Time Cleared" in df.columns:
        df["Duration (Minutes)"] = (
            (df["Time Cleared"] - df["Time Received"]).dt.total_seconds() / 60
        )
    else:
        df["Duration (Minutes)"] = None

    for col in ["Level", "BIU", "Device", "Name", "Message"]:
        if col in df.columns:
            df[col] = df[col].fillna("Unknown").astype(str).str.strip()

    return df


def fig_to_png_bytes(fig):
    return pio.to_image(fig, format="png", width=1400, height=800, scale=2)


def add_download_button(fig, file_name, label):
    try:
        img_bytes = fig_to_png_bytes(fig)
        st.download_button(
            label=label,
            data=img_bytes,
            file_name=file_name,
            mime="image/png",
            key=file_name
        )
    except Exception:
        st.info("PNG export ke liye kaleido install hona chahiye.")


# ----------------------------
# PPT helpers
# ----------------------------
def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    try:
        slide.placeholders[1].text = subtitle
    except Exception:
        pass


def update_template_second_slide_file_name(prs, file_name):
    if len(prs.slides) < 2:
        return

    slide = prs.slides[1]

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue

        try:
            text = shape.text_frame.text.strip()
            lower_text = text.lower()

            if "folder name" in lower_text or "foldername" in lower_text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = file_name
                p.font.size = Pt(20)
                p.font.bold = True
                return
        except Exception:
            pass


def add_summary_slide(prs, df, source_file, sheet_name):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"

    total_records = len(df)
    unique_levels = df["Level"].nunique() if "Level" in df.columns else 0
    unique_bius = df["BIU"].nunique() if "BIU" in df.columns else 0
    unique_devices = df["Device"].nunique() if "Device" in df.columns else 0
    unique_names = df["Name"].nunique() if "Name" in df.columns else 0

    avg_duration = (
        round(df["Duration (Minutes)"].dropna().mean(), 2)
        if "Duration (Minutes)" in df.columns and df["Duration (Minutes)"].notna().any()
        else 0
    )

    top_level = "N/A"
    if "Level" in df.columns and not df["Level"].dropna().empty:
        top_level = df["Level"].value_counts().idxmax()

    top_biu = "N/A"
    if "BIU" in df.columns and not df["BIU"].dropna().empty:
        top_biu = df["BIU"].value_counts().idxmax()

    top_device = "N/A"
    if "Device" in df.columns and not df["Device"].dropna().empty:
        top_device = df["Device"].value_counts().idxmax()

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    sections = [
        ("Dataset Overview", [
            f"Total alarm records: {total_records}",
            f"Unique alarm levels: {unique_levels}",
            f"Unique BIUs: {unique_bius}",
            f"Unique devices: {unique_devices}",
            f"Unique alarm names: {unique_names}"
        ]),
        ("Key Highlights", [
            f"Most frequent alarm level: {top_level}",
            f"Highest alarm BIU: {top_biu}",
            f"Most affected device: {top_device}",
            f"Average alarm duration: {avg_duration} minutes"
        ])
    ]

    first_written = False
    for heading, bullets in sections:
        if not first_written:
            p = body.paragraphs[0]
            first_written = True
        else:
            p = body.add_paragraph()

        p.text = heading
        p.level = 0
        p.font.bold = True
        p.font.size = Pt(20)

        for bullet in bullets:
            bp = body.add_paragraph()
            bp.text = bullet
            bp.level = 1
            bp.font.size = Pt(16)


def add_native_chart_slide(
    prs,
    title,
    categories,
    values,
    chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
    series_name="Count"
):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    chart_data.add_series(series_name, list(values))

    x = 700000
    y = 1200000
    cx = 7900000
    cy = 4700000

    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    if chart_type != XL_CHART_TYPE.PIE:
        try:
            chart.value_axis.has_major_gridlines = True
            chart.category_axis.tick_labels.font.size = Pt(12)
            chart.value_axis.tick_labels.font.size = Pt(12)
        except Exception:
            pass

    return slide


def add_table_slide(prs, title, df_table):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows, cols = df_table.shape[0] + 1, df_table.shape[1]
    table = slide.shapes.add_table(rows, cols, 350000, 1300000, 8900000, 4200000).table

    for j, col in enumerate(df_table.columns):
        table.cell(0, j).text = str(col)

    for i in range(df_table.shape[0]):
        for j in range(df_table.shape[1]):
            table.cell(i + 1, j).text = str(df_table.iloc[i, j])

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(10)


def add_template_summary_slide(prs, df, source_file, sheet_name):
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])

    total_records = len(df)
    unique_levels = df["Level"].nunique() if "Level" in df.columns else 0
    unique_bius = df["BIU"].nunique() if "BIU" in df.columns else 0
    unique_devices = df["Device"].nunique() if "Device" in df.columns else 0
    unique_names = df["Name"].nunique() if "Name" in df.columns else 0

    avg_duration = (
        round(df["Duration (Minutes)"].dropna().mean(), 2)
        if "Duration (Minutes)" in df.columns and df["Duration (Minutes)"].notna().any()
        else 0
    )

    top_level = "N/A"
    if "Level" in df.columns and not df["Level"].dropna().empty:
        top_level = df["Level"].value_counts().idxmax()

    top_biu = "N/A"
    if "BIU" in df.columns and not df["BIU"].dropna().empty:
        top_biu = df["BIU"].value_counts().idxmax()

    top_device = "N/A"
    if "Device" in df.columns and not df["Device"].dropna().empty:
        top_device = df["Device"].value_counts().idxmax()

    sections = [
        ("Dataset Overview", [
            f"Total alarm records: {total_records}",
            f"Unique alarm levels: {unique_levels}",
            f"Unique BIUs: {unique_bius}",
            f"Unique devices: {unique_devices}",
            f"Unique alarm names: {unique_names}"
        ]),
        ("Key Highlights", [
            f"Most frequent alarm level: {top_level}",
            f"Highest alarm BIU: {top_biu}",
            f"Most affected device: {top_device}",
            f"Average alarm duration: {avg_duration} minutes"
        ])
    ]

    try:
        if slide.shapes.title:
            slide.shapes.title.text = "Executive Summary"

        target_tf = None
        for shape in slide.placeholders:
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                target_tf = shape.text_frame
                break

        if target_tf is None:
            textbox = slide.shapes.add_textbox(600000, 1400000, 8400000, 3800000)
            target_tf = textbox.text_frame

        target_tf.clear()

        first_written = False
        for heading, bullets in sections:
            if not first_written:
                p = target_tf.paragraphs[0]
                first_written = True
            else:
                p = target_tf.add_paragraph()

            p.text = heading
            p.level = 0
            p.font.bold = True
            p.font.size = Pt(19)

            for bullet in bullets:
                bp = target_tf.add_paragraph()
                bp.text = bullet
                bp.level = 1
                bp.font.size = Pt(16)

    except Exception:
        pass


def build_editable_ppt_report(df, source_file, sheet_name, template_bytes=None):
    prs = Presentation(io.BytesIO(template_bytes)) if template_bytes else Presentation()

    if template_bytes:
        update_template_second_slide_file_name(prs, source_file)
        add_template_summary_slide(prs, df, source_file, sheet_name)
    else:
        add_title_slide(
            prs,
            "Alarm Analysis Report",
            f"File: {source_file}\nSheet: {sheet_name}\nGenerated from Streamlit Dashboard"
        )
        add_summary_slide(prs, df, source_file, sheet_name)

    if "Level" in df.columns and not df.empty:
        x = df["Level"].value_counts().reset_index()
        x.columns = ["Level", "Count"]
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Alarm Count by Level", x["Level"], x["Count"], chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED)

    if "Level" in df.columns and not df.empty:
        x = df["Level"].value_counts().reset_index()
        x.columns = ["Level", "Count"]
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Alarm Distribution by Level", x["Level"], x["Count"], chart_type=XL_CHART_TYPE.PIE)

    if "BIU" in df.columns and not df.empty:
        x = df["BIU"].value_counts().reset_index()
        x.columns = ["BIU", "Count"]
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Alarm Count by BIU", x["BIU"], x["Count"], chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED)

    if "Device" in df.columns and not df.empty:
        x = df["Device"].value_counts().reset_index().head(20)
        x.columns = ["Device", "Count"]
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Alarm Count by Device", x["Device"], x["Count"], chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED)

    if "Name" in df.columns and not df.empty:
        x = df["Name"].value_counts().reset_index().head(15)
        x.columns = ["Name", "Count"]
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Top Alarm Names", x["Name"], x["Count"], chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED)

    if "Time Received" in df.columns and df["Time Received"].notna().any():
        tdf = df.dropna(subset=["Time Received"]).copy()
        tdf["Date"] = tdf["Time Received"].dt.date.astype(str)
        x = tdf.groupby("Date").size().reset_index(name="Count")
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Alarm Timeline by Date", x["Date"], x["Count"], chart_type=XL_CHART_TYPE.LINE_MARKERS)

    if "Duration (Minutes)" in df.columns and df["Duration (Minutes)"].notna().any():
        dur = df["Duration (Minutes)"].dropna()
        bins = pd.cut(dur, bins=8)
        x = bins.value_counts().sort_index().reset_index()
        x.columns = ["Range", "Count"]
        x["Range"] = x["Range"].astype(str)
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Alarm Duration Distribution", x["Range"], x["Count"], chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED)

    if "Message" in df.columns and not df.empty:
        x = df["Message"].value_counts().reset_index().head(10)
        x.columns = ["Message", "Count"]
        x["Count"] = x["Count"].astype(int)
        add_native_chart_slide(prs, "Top 10 Messages", x["Message"], x["Count"], chart_type=XL_CHART_TYPE.BAR_CLUSTERED)

    if "Level" in df.columns and "BIU" in df.columns and not df.empty:
        heat_df = df.groupby(["Level", "BIU"]).size().reset_index(name="Count")
        heat_pivot = heat_df.pivot(index="Level", columns="BIU", values="Count").fillna(0).astype(int).reset_index()
        add_table_slide(prs, "Level vs BIU Matrix", heat_pivot)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# ----------------------------
# Chart config
# ----------------------------
plot_config = {
    "displaylogo": False,
    "toImageButtonOptions": {
        "format": "png",
        "filename": "chart_export",
        "height": 700,
        "width": 1200,
        "scale": 2
    }
}


# ----------------------------
# Upload area
# ----------------------------
uploaded_files = st.file_uploader(
    "Upload Excel Files",
    type=["xls", "xlsx"],
    accept_multiple_files=True
)

ppt_template = st.file_uploader(
    "Optional: Upload PowerPoint Template",
    type=["pptx"]
)

if uploaded_files:
    workbook_map = {}

    for uf in uploaded_files:
        try:
            workbook_map[uf.name] = {
                "file": uf,
                "sheets": load_workbook_sheets(uf)
            }
        except Exception as e:
            st.error(f"Could not read workbook '{uf.name}': {e}")

    valid_workbooks = list(workbook_map.keys())

    if valid_workbooks:
        st.sidebar.header("Workbook / Sheet Selection")

        selected_workbook = st.sidebar.selectbox("Select Workbook", valid_workbooks)
        selected_sheet = st.sidebar.selectbox(
            "Select Sheet",
            workbook_map[selected_workbook]["sheets"]
        )

        try:
            raw_df = load_selected_sheet(workbook_map[selected_workbook]["file"], selected_sheet)
            df = clean_dataframe(raw_df)

            st.success(f"Loaded workbook: {selected_workbook} | Sheet: {selected_sheet}")

            st.subheader("Data Preview")
            st.dataframe(df, use_container_width=True)

            # Filters
            st.sidebar.header("Filters")
            filtered_df = df.copy()

            if "Level" in filtered_df.columns:
                level_options = sorted(filtered_df["Level"].dropna().unique().tolist())
                selected_levels = st.sidebar.multiselect("Select Level", level_options, default=level_options)
                filtered_df = filtered_df[filtered_df["Level"].isin(selected_levels)]

            if "BIU" in filtered_df.columns:
                biu_options = sorted(filtered_df["BIU"].dropna().unique().tolist())
                selected_bius = st.sidebar.multiselect("Select BIU", biu_options, default=biu_options)
                filtered_df = filtered_df[filtered_df["BIU"].isin(selected_bius)]

            if "Device" in filtered_df.columns:
                device_options = sorted(filtered_df["Device"].dropna().unique().tolist())
                selected_devices = st.sidebar.multiselect("Select Device", device_options, default=device_options)
                filtered_df = filtered_df[filtered_df["Device"].isin(selected_devices)]

            if "Name" in filtered_df.columns:
                name_options = sorted(filtered_df["Name"].dropna().unique().tolist())
                selected_names = st.sidebar.multiselect("Select Name", name_options, default=name_options)
                filtered_df = filtered_df[filtered_df["Name"].isin(selected_names)]

            # Summary
            st.subheader("Summary")
            c1, c2, c3, c4 = st.columns(4)

            total_records = len(filtered_df)
            unique_levels = filtered_df["Level"].nunique() if "Level" in filtered_df.columns else 0
            unique_bius = filtered_df["BIU"].nunique() if "BIU" in filtered_df.columns else 0
            avg_duration = (
                round(filtered_df["Duration (Minutes)"].dropna().mean(), 2)
                if "Duration (Minutes)" in filtered_df.columns and filtered_df["Duration (Minutes)"].notna().any()
                else 0
            )

            c1.metric("Total Records", total_records)
            c2.metric("Unique Levels", unique_levels)
            c3.metric("Unique BIUs", unique_bius)
            c4.metric("Avg Duration (Min)", avg_duration)

            # Row 1
            col1, col2 = st.columns(2)

            if "Level" in filtered_df.columns and not filtered_df.empty:
                level_count = filtered_df["Level"].value_counts().reset_index()
                level_count.columns = ["Level", "Count"]
                level_count["Count"] = level_count["Count"].astype(int)

                fig1 = px.bar(level_count, x="Level", y="Count", text="Count", title="Alarm Count by Level")
                fig1.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                fig1.update_layout(xaxis_title="Level", yaxis_title="Count", yaxis=dict(tickformat="d"))
                fig1 = apply_theme(fig1)
                col1.plotly_chart(fig1, use_container_width=True, config=plot_config)
                add_download_button(fig1, "alarm_count_by_level.png", "Download Level Chart")

            if "BIU" in filtered_df.columns and not filtered_df.empty:
                biu_count = filtered_df["BIU"].value_counts().reset_index()
                biu_count.columns = ["BIU", "Count"]
                biu_count["Count"] = biu_count["Count"].astype(int)

                fig2 = px.bar(biu_count, x="BIU", y="Count", text="Count", title="Alarm Count by BIU")
                fig2.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                fig2.update_layout(xaxis_title="BIU", yaxis_title="Count", yaxis=dict(tickformat="d"))
                fig2 = apply_theme(fig2)
                col2.plotly_chart(fig2, use_container_width=True, config=plot_config)
                add_download_button(fig2, "alarm_count_by_biu.png", "Download BIU Chart")

            # Pie chart
            if "Level" in filtered_df.columns and not filtered_df.empty:
                st.subheader("Level Distribution Pie Chart")
                pie_df = filtered_df["Level"].value_counts().reset_index()
                pie_df.columns = ["Level", "Count"]
                pie_df["Count"] = pie_df["Count"].astype(int)

                fig_pie = px.pie(
                    pie_df,
                    names="Level",
                    values="Count",
                    title="Alarm Distribution by Level",
                    hole=0.0
                )
                fig_pie = apply_theme(fig_pie)
                st.plotly_chart(fig_pie, use_container_width=True, config=plot_config)
                add_download_button(fig_pie, "alarm_distribution_by_level_pie.png", "Download Pie Chart")

            # Row 2
            col3, col4 = st.columns(2)

            if "Device" in filtered_df.columns and not filtered_df.empty:
                device_count = filtered_df["Device"].value_counts().reset_index().head(20)
                device_count.columns = ["Device", "Count"]
                device_count["Count"] = device_count["Count"].astype(int)

                fig3 = px.bar(device_count, x="Device", y="Count", text="Count", title="Alarm Count by Device")
                fig3.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                fig3.update_layout(xaxis_title="Device", yaxis_title="Count", yaxis=dict(tickformat="d"))
                fig3 = apply_theme(fig3)
                col3.plotly_chart(fig3, use_container_width=True, config=plot_config)
                add_download_button(fig3, "alarm_count_by_device.png", "Download Device Chart")

            if "Name" in filtered_df.columns and not filtered_df.empty:
                name_count = filtered_df["Name"].value_counts().reset_index().head(15)
                name_count.columns = ["Name", "Count"]
                name_count["Count"] = name_count["Count"].astype(int)

                fig4 = px.bar(name_count, x="Name", y="Count", text="Count", title="Top Alarm Names")
                fig4.update_traces(textposition="outside", texttemplate="%{text:.0f}")
                fig4.update_layout(xaxis_title="Name", yaxis_title="Count", yaxis=dict(tickformat="d"))
                fig4 = apply_theme(fig4)
                col4.plotly_chart(fig4, use_container_width=True, config=plot_config)
                add_download_button(fig4, "top_alarm_names.png", "Download Name Chart")

            # Timeline
            if "Time Received" in filtered_df.columns and filtered_df["Time Received"].notna().any():
                st.subheader("Alarm Timeline")
                timeline_df = filtered_df.dropna(subset=["Time Received"]).copy()
                timeline_df["Date"] = timeline_df["Time Received"].dt.date
                timeline_count = timeline_df.groupby("Date").size().reset_index(name="Count")
                timeline_count["Count"] = timeline_count["Count"].astype(int)

                fig5 = px.line(
                    timeline_count,
                    x="Date",
                    y="Count",
                    markers=True,
                    title="Alarm Timeline by Date"
                )
                fig5.update_traces(
                    hovertemplate="Date: %{x}<br>Alarm Count: %{y:d}<extra></extra>"
                )
                fig5.update_layout(
                    xaxis_title="Date",
                    yaxis_title="Alarm Count",
                    yaxis=dict(tickformat="d", dtick=1)
                )
                fig5 = apply_theme(fig5)
                st.plotly_chart(fig5, use_container_width=True, config=plot_config)
                add_download_button(fig5, "alarm_timeline.png", "Download Timeline Chart")

            # Duration
            if "Duration (Minutes)" in filtered_df.columns and filtered_df["Duration (Minutes)"].notna().any():
                st.subheader("Alarm Duration Analysis")
                fig6 = px.histogram(
                    filtered_df.dropna(subset=["Duration (Minutes)"]),
                    x="Duration (Minutes)",
                    nbins=30,
                    title="Alarm Duration Distribution"
                )
                fig6.update_layout(xaxis_title="Duration (Minutes)", yaxis_title="Frequency")
                fig6 = apply_theme(fig6)
                st.plotly_chart(fig6, use_container_width=True, config=plot_config)
                add_download_button(fig6, "duration_distribution.png", "Download Duration Chart")

            # Heatmap
            if "Level" in filtered_df.columns and "BIU" in filtered_df.columns and not filtered_df.empty:
                st.subheader("Heatmap View")
                heat_df = filtered_df.groupby(["Level", "BIU"]).size().reset_index(name="Count")
                heat_pivot = heat_df.pivot(index="Level", columns="BIU", values="Count").fillna(0)
                heat_pivot = heat_pivot.astype(int)

                max_val = int(heat_pivot.values.max()) if heat_pivot.values.size > 0 else 1
                if max_val == 0:
                    max_val = 1

                fig8 = go.Figure(
                    data=go.Heatmap(
                        z=heat_pivot.values,
                        x=heat_pivot.columns,
                        y=heat_pivot.index,
                        text=heat_pivot.values,
                        texttemplate="%{text:d}",
                        colorscale=[
                            [0.0, "#fffde7"],
                            [0.25, "#fff176"],
                            [0.50, "#66bb6a"],
                            [0.75, "#43a047"],
                            [1.0, "#d32f2f"]
                        ],
                        zmin=0,
                        zmax=max_val,
                        colorbar=dict(title="Count", tickformat="d"),
                        hovertemplate="Level: %{y}<br>BIU: %{x}<br>Count: %{z:d}<extra></extra>"
                    )
                )
                fig8.update_layout(
                    title="Heatmap: Level vs BIU",
                    xaxis_title="BIU",
                    yaxis_title="Level"
                )
                fig8 = apply_theme(fig8)
                st.plotly_chart(fig8, use_container_width=True, config=plot_config)
                add_download_button(fig8, "heatmap_level_vs_biu.png", "Download Heatmap")

            # Sunburst
            hierarchy_cols = [c for c in ["Level", "BIU", "Device", "Name"] if c in filtered_df.columns]
            if len(hierarchy_cols) >= 2 and not filtered_df.empty:
                st.subheader("Relationship View: Level → BIU → Device → Name")
                sun_df = filtered_df.groupby(hierarchy_cols).size().reset_index(name="Count")
                sun_df["Count"] = sun_df["Count"].astype(int)

                fig_sunburst = px.sunburst(
                    sun_df,
                    path=hierarchy_cols,
                    values="Count",
                    title="Level → BIU → Device → Name"
                )
                fig_sunburst = apply_theme(fig_sunburst)
                st.plotly_chart(fig_sunburst, use_container_width=True, config=plot_config)
                add_download_button(
                    fig_sunburst,
                    "level_biu_device_name_sunburst.png",
                    "Download Sunburst Chart"
                )

            # Top messages
            if "Message" in filtered_df.columns and not filtered_df.empty:
                st.subheader("Top Messages")
                message_count = filtered_df["Message"].value_counts().reset_index().head(10)
                message_count.columns = ["Message", "Count"]
                message_count["Count"] = message_count["Count"].astype(int)

                fig9 = px.bar(
                    message_count,
                    x="Count",
                    y="Message",
                    orientation="h",
                    text="Count",
                    title="Top 10 Messages"
                )
                fig9.update_traces(texttemplate="%{text:.0f}")
                fig9.update_layout(
                    yaxis={"categoryorder": "total ascending"},
                    xaxis=dict(tickformat="d")
                )
                fig9 = apply_theme(fig9)
                st.plotly_chart(fig9, use_container_width=True, config=plot_config)
                add_download_button(fig9, "top_messages.png", "Download Message Chart")

            # PPT
            st.subheader("PowerPoint Report")
            if not filtered_df.empty:
                try:
                    template_bytes = ppt_template.getvalue() if ppt_template is not None else None

                    ppt_buffer = build_editable_ppt_report(
                        filtered_df,
                        selected_workbook,
                        selected_sheet,
                        template_bytes=template_bytes
                    )

                    out_name = f"Alarm_Report_Editable_{selected_workbook.replace('.', '_')}_{selected_sheet}.pptx"

                    st.download_button(
                        label="Download Editable PowerPoint Report",
                        data=ppt_buffer,
                        file_name=out_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

                    st.success("Editable PPT report ready. Is file me charts PowerPoint ke andar edit ho sakte hain.")

                except Exception as e:
                    st.error(f"PPT generation failed: {e}")

        except Exception as e:
            st.error(f"Could not load selected sheet: {e}")
else:
    st.info("Please upload one or more Excel files to continue.")